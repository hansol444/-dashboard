"""
create.py -- 장표 제작 에이전트

8단계 워크플로를 통해 컨텍스트/입력 텍스트로부터
구조화된 프레젠테이션(PPTX)을 자동 생성한다.

사용법:
  # Single-shot (확인 단계 생략, 자동 생성)
  python create.py --input "트랜스크립트 또는 컨텍스트" --output output/result.pptx

  # Session mode (단계별 진행)
  python create.py --session abc123 --step 1 --input "initial context"
  python create.py --session abc123 --step 2 --input "답변: 경영진 대상, 전략 브리핑"
  python create.py --session abc123 --step 3
  python create.py --session abc123 --step 4 --input "확정"
  python create.py --session abc123 --step 5
  python create.py --session abc123 --step 6 --input "확정"
  python create.py --session abc123 --step 7
"""

import sys
import os
import argparse
import json
import logging
import uuid
from pathlib import Path

import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# .env 파일 자동 로드
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

# ─── 설정 ──────────────────────────────────────────────────────────────────────

ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
MODEL = "claude-sonnet-4-6"
SESSIONS_DIR = Path(__file__).parent / "sessions"
OUTPUT_DIR = Path(__file__).parent / "output"

# 슬라이드 디자인 상수
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
HEADER_HEIGHT = Inches(0.8)
FOOTER_HEIGHT = Inches(0.4)
FONT_NAME = "맑은 고딕"
COLOR_HEADER_BG = RGBColor(0x1B, 0x36, 0x5D)   # dark navy
COLOR_ACCENT = RGBColor(0x00, 0xB3, 0x88)       # teal green
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)
COLOR_DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
COLOR_LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

LAYOUT_TYPES = ["텍스트형", "표형", "차트형", "다이어그램형", "타임라인형", "하이브리드형"]

logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
logger = logging.getLogger(__name__)

# ─── API 클라이언트 ────────────────────────────────────────────────────────────

def get_api_client() -> anthropic.Anthropic:
    api_key = ANTHROPIC_API_KEY
    if not api_key:
        raise ValueError(
            "ANTHROPIC_API_KEY가 설정되지 않았습니다.\n"
            "환경변수로 설정하거나 .env 파일에 ANTHROPIC_API_KEY=sk-... 형태로 추가하세요."
        )
    return anthropic.Anthropic(api_key=api_key)


def call_claude(system_prompt: str, user_prompt: str) -> str:
    """Claude API를 호출하고 텍스트 응답을 반환한다."""
    client = get_api_client()
    message = client.messages.create(
        model=MODEL,
        max_tokens=4096,
        system=system_prompt,
        messages=[{"role": "user", "content": user_prompt}],
    )
    return message.content[0].text


def call_claude_json(system_prompt: str, user_prompt: str) -> dict:
    """Claude API를 호출하고 JSON 파싱된 응답을 반환한다."""
    raw = call_claude(system_prompt, user_prompt)
    # JSON 블록 추출 (```json ... ``` 또는 순수 JSON)
    if "```json" in raw:
        raw = raw.split("```json", 1)[1].split("```", 1)[0]
    elif "```" in raw:
        raw = raw.split("```", 1)[1].split("```", 1)[0]
    return json.loads(raw.strip())


# ─── 세션 관리 ─────────────────────────────────────────────────────────────────

def new_session(session_id: str = None) -> dict:
    """새 세션 상태를 생성한다."""
    sid = session_id or uuid.uuid4().hex[:12]
    return {
        "id": sid,
        "step": 0,
        "context": "",
        "audience": "",
        "purpose": "",
        "headMessages": [],
        "layouts": [],
        "slides": [],
        "outputPath": "",
    }


def load_session(session_id: str) -> dict:
    SESSIONS_DIR.mkdir(parents=True, exist_ok=True)
    path = SESSIONS_DIR / f"{session_id}.json"
    if path.exists():
        return json.loads(path.read_text(encoding="utf-8"))
    return new_session(session_id)


def save_session(session: dict):
    SESSIONS_DIR.mkdir(parents=True, exist_ok=True)
    path = SESSIONS_DIR / f"{session['id']}.json"
    path.write_text(json.dumps(session, ensure_ascii=False, indent=2), encoding="utf-8")


# ─── 응답 포맷 ─────────────────────────────────────────────────────────────────

def make_response(step: int, status: str, message: str, data: dict,
                  next_step: int, needs_input: bool, prompt: str = "") -> dict:
    return {
        "step": step,
        "status": status,
        "message": message,
        "data": data,
        "nextStep": next_step,
        "needsInput": needs_input,
        "prompt": prompt,
    }


# ─── Step 1: 입력 분석 ────────────────────────────────────────────────────────

def step1_input_analysis(session: dict, user_input: str) -> dict:
    """입력 텍스트를 분석하여 유형을 감지하고 핵심 내용을 추출한다."""
    # 파일 경로인 경우 읽기
    if os.path.isfile(user_input):
        content = Path(user_input).read_text(encoding="utf-8")
    else:
        content = user_input

    session["context"] = content

    system_prompt = (
        "당신은 프레젠테이션 제작 전문가입니다. "
        "사용자가 제공한 텍스트의 유형을 분석하세요.\n"
        "반드시 JSON만 반환하세요. 다른 텍스트는 포함하지 마세요.\n"
        "JSON 형식:\n"
        '{"inputType": "transcript|draft|conversation|document|bullet_points", '
        '"summary": "핵심 내용 요약 (3문장 이내)", '
        '"keyTopics": ["주제1", "주제2", ...], '
        '"estimatedSlides": 숫자, '
        '"language": "ko|en"}'
    )

    result = call_claude_json(system_prompt, f"다음 텍스트를 분석하세요:\n\n{content[:3000]}")
    session["step"] = 1
    save_session(session)

    return make_response(
        step=1,
        status="done",
        message="입력을 분석했습니다.",
        data=result,
        next_step=2,
        needs_input=True,
        prompt=(
            "입력이 분석되었습니다. 다음 질문에 답변해 주세요:\n"
            "1) 누구 입장에서 발표하나요? (대상 청중)\n"
            "2) 어느 깊이까지 다루나요? (개요/상세)\n"
            "3) 현재 상태는? (초안/최종)"
        ),
    )


# ─── Step 2: 명확화 질문 ──────────────────────────────────────────────────────

def step2_clarification(session: dict, user_input: str) -> dict:
    """사용자 답변을 반영하고 3가지 명확화 질문을 정리한다."""
    system_prompt = (
        "당신은 프레젠테이션 기획자입니다. "
        "사용자의 원본 컨텍스트와 답변을 바탕으로 청중, 목적, 깊이를 정리하세요.\n"
        "반드시 JSON만 반환하세요.\n"
        "JSON 형식:\n"
        '{"audience": "청중 설명", '
        '"purpose": "발표 목적", '
        '"depth": "개요|상세|전략 브리핑", '
        '"suggestedSlideCount": 숫자, '
        '"slideTopics": ["슬라이드1 주제", "슬라이드2 주제", ...]}'
    )

    user_msg = (
        f"원본 컨텍스트:\n{session['context'][:2000]}\n\n"
        f"사용자 답변:\n{user_input}"
    )

    result = call_claude_json(system_prompt, user_msg)
    session["audience"] = result.get("audience", "")
    session["purpose"] = result.get("purpose", "")
    session["step"] = 2
    save_session(session)

    return make_response(
        step=2,
        status="done",
        message="명확화가 완료되었습니다.",
        data=result,
        next_step=3,
        needs_input=False,
        prompt="다음 단계에서 헤드메시지 초안을 생성합니다.",
    )


# ─── Step 3: 헤드메시지 초안 ──────────────────────────────────────────────────

def step3_head_messages(session: dict) -> dict:
    """슬라이드별 헤드메시지를 생성하고 9-체크리스트로 검증한다."""
    system_prompt = (
        "당신은 프레젠테이션 헤드메시지 전문가입니다.\n"
        "주어진 컨텍스트, 청중, 목적을 기반으로 슬라이드별 헤드메시지를 작성하세요.\n\n"
        "9-체크리스트 검증 기준:\n"
        "1. 한 문장으로 핵심 메시지 전달\n"
        "2. So What 테스트 통과 (행동/판단 유도)\n"
        "3. 수치/근거 포함 여부\n"
        "4. 청중 관점에서 의미 있는 내용\n"
        "5. 이전 슬라이드와 논리적 연결\n"
        "6. 중복 없음\n"
        "7. 추상적 표현 배제\n"
        "8. 15단어 이내 간결성\n"
        "9. 전체 스토리라인에서 역할 명확\n\n"
        "반드시 JSON만 반환하세요.\n"
        "JSON 형식:\n"
        '{"slides": [\n'
        '  {"slideNumber": 1, "headMessage": "...", "subPoints": ["...", "..."], '
        '"checklist": {"soWhat": true, "hasEvidence": true, ...}, "role": "도입|본론|결론"},\n'
        '  ...\n'
        ']}'
    )

    user_msg = (
        f"컨텍스트:\n{session['context'][:2000]}\n\n"
        f"청중: {session['audience']}\n"
        f"목적: {session['purpose']}"
    )

    result = call_claude_json(system_prompt, user_msg)
    session["headMessages"] = result.get("slides", [])
    session["step"] = 3
    save_session(session)

    return make_response(
        step=3,
        status="done",
        message="헤드메시지 초안을 생성했습니다.",
        data=result,
        next_step=4,
        needs_input=True,
        prompt="초안을 검토하고 수정사항이 있으면 알려주세요. 없으면 '확정'이라고 해주세요.",
    )


# ─── Step 4: 사용자 확인 ──────────────────────────────────────────────────────

def step4_user_confirmation(session: dict, user_input: str) -> dict:
    """사용자 확정 또는 수정 요청을 처리한다."""
    if "확정" in user_input:
        session["step"] = 4
        save_session(session)
        return make_response(
            step=4,
            status="done",
            message="헤드메시지가 확정되었습니다.",
            data={"confirmed": True, "slides": session["headMessages"]},
            next_step=5,
            needs_input=False,
            prompt="다음 단계에서 레이아웃을 자동 선택합니다.",
        )

    # 수정 요청 처리
    system_prompt = (
        "당신은 프레젠테이션 헤드메시지 수정 전문가입니다.\n"
        "현재 헤드메시지 초안과 사용자의 수정 요청을 반영하여 수정된 버전을 반환하세요.\n"
        "반드시 JSON만 반환하세요.\n"
        "JSON 형식:\n"
        '{"slides": [\n'
        '  {"slideNumber": 1, "headMessage": "...", "subPoints": ["...", "..."], "role": "도입|본론|결론"},\n'
        '  ...\n'
        ']}'
    )

    user_msg = (
        f"현재 초안:\n{json.dumps(session['headMessages'], ensure_ascii=False)}\n\n"
        f"수정 요청:\n{user_input}"
    )

    result = call_claude_json(system_prompt, user_msg)
    session["headMessages"] = result.get("slides", [])
    session["step"] = 3  # stay at step 3-4 loop
    save_session(session)

    return make_response(
        step=4,
        status="revised",
        message="수정된 헤드메시지를 반영했습니다.",
        data=result,
        next_step=4,
        needs_input=True,
        prompt="수정된 초안을 확인하세요. 추가 수정이 필요하면 알려주세요. 없으면 '확정'이라고 해주세요.",
    )


# ─── Step 5: 레이아웃 선택 ────────────────────────────────────────────────────

def step5_layout_selection(session: dict) -> dict:
    """슬라이드별 최적 레이아웃을 자동 선택한다."""
    system_prompt = (
        "당신은 프레젠테이션 레이아웃 전문가입니다.\n"
        "각 슬라이드의 헤드메시지와 내용을 기반으로 최적의 레이아웃 유형을 선택하세요.\n\n"
        "레이아웃 유형:\n"
        '- "텍스트형": 텍스트 중심 슬라이드 (아젠다, 핵심 메시지)\n'
        '- "표형": 표/비교 슬라이드\n'
        '- "차트형": 차트/그래프 슬라이드 (bar, line, pie)\n'
        '- "다이어그램형": 프로세스/플로우 다이어그램\n'
        '- "타임라인형": 타임라인 슬라이드\n'
        '- "하이브리드형": 복합 레이아웃 (텍스트 + 차트)\n\n'
        "반드시 JSON만 반환하세요.\n"
        "JSON 형식:\n"
        '{"layouts": [\n'
        '  {"slideNumber": 1, "layout": "텍스트형", "reason": "이유 설명", '
        '"contentStructure": {"type": "bullets|table|chart|diagram|timeline|mixed", "details": "..."}},\n'
        '  ...\n'
        ']}'
    )

    user_msg = (
        f"슬라이드 목록:\n{json.dumps(session['headMessages'], ensure_ascii=False)}"
    )

    result = call_claude_json(system_prompt, user_msg)
    session["layouts"] = result.get("layouts", [])
    session["step"] = 5
    save_session(session)

    return make_response(
        step=5,
        status="done",
        message="레이아웃을 자동 선택했습니다.",
        data=result,
        next_step=6,
        needs_input=True,
        prompt="레이아웃을 확인하세요. 변경이 필요하면 알려주세요. 없으면 '확정'이라고 해주세요.",
    )


# ─── Step 6: 레이아웃 확인 ────────────────────────────────────────────────────

def step6_layout_confirmation(session: dict, user_input: str) -> dict:
    """레이아웃 확정 또는 수정 요청을 처리한다."""
    if "확정" in user_input:
        session["step"] = 6
        save_session(session)
        return make_response(
            step=6,
            status="done",
            message="레이아웃이 확정되었습니다.",
            data={"confirmed": True, "layouts": session["layouts"]},
            next_step=7,
            needs_input=False,
            prompt="다음 단계에서 PPTX를 생성합니다.",
        )

    # 수정 요청 처리
    system_prompt = (
        "당신은 프레젠테이션 레이아웃 수정 전문가입니다.\n"
        "현재 레이아웃 설정과 사용자의 수정 요청을 반영하여 수정된 버전을 반환하세요.\n"
        "레이아웃 유형: 텍스트형, 표형, 차트형, 다이어그램형, 타임라인형, 하이브리드형\n"
        "반드시 JSON만 반환하세요.\n"
        "JSON 형식:\n"
        '{"layouts": [\n'
        '  {"slideNumber": 1, "layout": "텍스트형", "reason": "..."},\n'
        '  ...\n'
        ']}'
    )

    user_msg = (
        f"현재 레이아웃:\n{json.dumps(session['layouts'], ensure_ascii=False)}\n\n"
        f"수정 요청:\n{user_input}"
    )

    result = call_claude_json(system_prompt, user_msg)
    session["layouts"] = result.get("layouts", [])
    session["step"] = 5  # stay at step 5-6 loop
    save_session(session)

    return make_response(
        step=6,
        status="revised",
        message="수정된 레이아웃을 반영했습니다.",
        data=result,
        next_step=6,
        needs_input=True,
        prompt="수정된 레이아웃을 확인하세요. 추가 수정이 필요하면 알려주세요. 없으면 '확정'이라고 해주세요.",
    )


# ─── Step 7: PPTX 생성 ────────────────────────────────────────────────────────

def _generate_slide_content(session: dict) -> list[dict]:
    """Claude를 사용하여 각 슬라이드의 상세 콘텐츠를 생성한다."""
    system_prompt = (
        "당신은 프레젠테이션 콘텐츠 작성 전문가입니다.\n"
        "각 슬라이드의 헤드메시지, 레이아웃 유형, 원본 컨텍스트를 기반으로 "
        "슬라이드에 들어갈 구체적인 콘텐츠를 생성하세요.\n\n"
        "레이아웃별 콘텐츠 형식:\n"
        "- 텍스트형: bullets 배열 (각 항목은 문자열)\n"
        "- 표형: table 객체 {headers: [...], rows: [[...], [...]]}\n"
        "- 차트형: chart 객체 {chartType: 'bar|line|pie', labels: [...], values: [...], title: '...'}\n"
        "- 다이어그램형: steps 배열 [{label: '...', description: '...'}]\n"
        "- 타임라인형: events 배열 [{date: '...', title: '...', description: '...'}]\n"
        "- 하이브리드형: {text: [...], chart: {chartType: '...', labels: [...], values: [...]}}\n\n"
        "반드시 JSON만 반환하세요.\n"
        "JSON 형식:\n"
        '{"slides": [\n'
        '  {"slideNumber": 1, "title": "...", "layout": "텍스트형", "content": {...}},\n'
        '  ...\n'
        ']}'
    )

    slides_info = []
    for i, hm in enumerate(session["headMessages"]):
        layout_info = session["layouts"][i] if i < len(session["layouts"]) else {"layout": "텍스트형"}
        slides_info.append({
            "slideNumber": hm.get("slideNumber", i + 1),
            "headMessage": hm.get("headMessage", ""),
            "subPoints": hm.get("subPoints", []),
            "layout": layout_info.get("layout", "텍스트형"),
        })

    user_msg = (
        f"컨텍스트:\n{session['context'][:2000]}\n\n"
        f"청중: {session['audience']}\n"
        f"목적: {session['purpose']}\n\n"
        f"슬라이드 구성:\n{json.dumps(slides_info, ensure_ascii=False)}"
    )

    result = call_claude_json(system_prompt, user_msg)
    return result.get("slides", [])


def _add_header(slide, title: str):
    """슬라이드 상단에 헤더 영역을 추가한다."""
    header_shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0),
        SLIDE_WIDTH, HEADER_HEIGHT,
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_HEADER_BG
    header_shape.line.fill.background()

    tf = header_shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.LEFT


def _add_footer(slide, page_number: int, company_name: str = "전략추진실"):
    """슬라이드 하단에 회사명 + 페이지 번호를 추가한다."""
    footer_left = Inches(0.5)
    footer_top = SLIDE_HEIGHT - FOOTER_HEIGHT
    footer_width = SLIDE_WIDTH - Inches(1.0)

    txBox = slide.shapes.add_textbox(footer_left, footer_top, footer_width, FOOTER_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{company_name}"
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_DARK_GRAY
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.LEFT

    # 페이지 번호 (오른쪽)
    page_box = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(1.0), footer_top, Inches(0.5), FOOTER_HEIGHT,
    )
    tf2 = page_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = str(page_number)
    p2.font.size = Pt(9)
    p2.font.color.rgb = COLOR_DARK_GRAY
    p2.font.name = FONT_NAME
    p2.alignment = PP_ALIGN.RIGHT


def _set_font(run, size_pt=14, bold=False, color=None, name=None):
    """run 객체의 폰트를 설정한다."""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = name or FONT_NAME
    if color:
        run.font.color.rgb = color


def _build_text_slide(slide, slide_data: dict):
    """텍스트형 슬라이드: 불릿 포인트 목록."""
    content = slide_data.get("content", {})
    bullets = content if isinstance(content, list) else content.get("bullets", [])
    if not bullets:
        bullets = slide_data.get("subPoints", ["(내용 없음)"])

    body_top = HEADER_HEIGHT + Inches(0.3)
    body_height = SLIDE_HEIGHT - HEADER_HEIGHT - FOOTER_HEIGHT - Inches(0.6)

    txBox = slide.shapes.add_textbox(Inches(0.8), body_top, SLIDE_WIDTH - Inches(1.6), body_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if isinstance(bullet, dict):
            bullet = bullet.get("text", str(bullet))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(4)
        p.level = 0
        run = p.add_run()
        run.text = f"  {bullet}"
        _set_font(run, size_pt=14, color=COLOR_DARK_GRAY)

        # 불릿 마커
        p.alignment = PP_ALIGN.LEFT


def _build_table_slide(slide, slide_data: dict):
    """표형 슬라이드: 표 생성."""
    content = slide_data.get("content", {})
    table_data = content if isinstance(content, dict) and "headers" in content else content.get("table", {})
    headers = table_data.get("headers", ["항목", "내용"])
    rows_data = table_data.get("rows", [["데이터 없음", "-"]])

    body_top = HEADER_HEIGHT + Inches(0.4)
    num_rows = len(rows_data) + 1  # +1 for header
    num_cols = len(headers)
    table_width = SLIDE_WIDTH - Inches(1.6)
    table_height = min(Inches(0.4) * num_rows, SLIDE_HEIGHT - body_top - FOOTER_HEIGHT - Inches(0.3))

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.8), body_top,
        table_width, table_height,
    )
    table = table_shape.table

    # 헤더 행
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(header)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.name = FONT_NAME
            paragraph.font.color.rgb = COLOR_WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_BG

    # 데이터 행
    for i, row in enumerate(rows_data):
        for j in range(num_cols):
            cell_text = str(row[j]) if j < len(row) else ""
            cell = table.cell(i + 1, j)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.name = FONT_NAME
                paragraph.font.color.rgb = COLOR_DARK_GRAY
                paragraph.alignment = PP_ALIGN.LEFT
            # 줄무늬 배경
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY


def _build_chart_slide(slide, slide_data: dict):
    """차트형 슬라이드: 차트 데이터를 텍스트+도형으로 시각화한다.

    python-pptx의 차트 기능 대신, 간단한 막대 시각화를 도형으로 표현한다.
    """
    content = slide_data.get("content", {})
    chart_data = content if isinstance(content, dict) and "labels" in content else content.get("chart", {})
    labels = chart_data.get("labels", ["A", "B", "C"])
    values = chart_data.get("values", [30, 50, 20])
    chart_title = chart_data.get("title", "")

    body_top = HEADER_HEIGHT + Inches(0.4)
    body_left = Inches(1.0)
    chart_area_width = SLIDE_WIDTH - Inches(2.0)
    chart_area_height = SLIDE_HEIGHT - body_top - FOOTER_HEIGHT - Inches(0.6)

    # 차트 제목
    if chart_title:
        title_box = slide.shapes.add_textbox(body_left, body_top, chart_area_width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = chart_title
        _set_font(run, size_pt=14, bold=True, color=COLOR_HEADER_BG)
        p.alignment = PP_ALIGN.CENTER
        body_top += Inches(0.5)
        chart_area_height -= Inches(0.5)

    # 수평 막대 차트를 도형으로 그리기
    max_val = max(values) if values else 1
    bar_height = min(Inches(0.5), chart_area_height / max(len(labels), 1) - Inches(0.1))
    max_bar_width = chart_area_width - Inches(2.0)  # 라벨 공간 확보

    for i, (label, value) in enumerate(zip(labels, values)):
        y_pos = body_top + Inches(0.1) + (bar_height + Inches(0.15)) * i

        # 라벨
        label_box = slide.shapes.add_textbox(body_left, y_pos, Inches(1.8), bar_height)
        tf = label_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(label)
        _set_font(run, size_pt=10, color=COLOR_DARK_GRAY)
        p.alignment = PP_ALIGN.RIGHT

        # 막대
        bar_width = max(Inches(0.2), int(max_bar_width * (value / max_val)))
        bar = slide.shapes.add_shape(
            1,  # RECTANGLE
            body_left + Inches(2.0), y_pos,
            bar_width, bar_height,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_ACCENT
        bar.line.fill.background()

        # 값 표시
        val_box = slide.shapes.add_textbox(
            body_left + Inches(2.0) + bar_width + Inches(0.1), y_pos,
            Inches(0.8), bar_height,
        )
        tf2 = val_box.text_frame
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = str(value)
        _set_font(run2, size_pt=10, bold=True, color=COLOR_HEADER_BG)


def _build_diagram_slide(slide, slide_data: dict):
    """다이어그램형 슬라이드: 프로세스 흐름을 단계별 도형으로 표현."""
    content = slide_data.get("content", {})
    steps = content if isinstance(content, list) else content.get("steps", [])
    if not steps:
        steps = [{"label": "단계 1", "description": "(내용 없음)"}]

    body_top = HEADER_HEIGHT + Inches(0.5)
    num_steps = len(steps)
    available_width = SLIDE_WIDTH - Inches(1.6)
    step_width = min(Inches(2.2), available_width / max(num_steps, 1) - Inches(0.2))
    gap = Inches(0.3) if num_steps > 1 else 0
    total_width = step_width * num_steps + gap * (num_steps - 1)
    start_x = (SLIDE_WIDTH - total_width) / 2

    for i, step_info in enumerate(steps):
        if isinstance(step_info, str):
            step_info = {"label": step_info, "description": ""}
        label = step_info.get("label", f"단계 {i+1}")
        desc = step_info.get("description", "")

        x_pos = start_x + (step_width + gap) * i
        y_pos = body_top + Inches(1.0)

        # 번호 원형
        circle = slide.shapes.add_shape(
            9,  # OVAL
            x_pos + (step_width - Inches(0.6)) / 2, body_top,
            Inches(0.6), Inches(0.6),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_ACCENT
        circle.line.fill.background()
        tf_c = circle.text_frame
        tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_c = tf_c.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        run_c = p_c.add_run()
        run_c.text = str(i + 1)
        _set_font(run_c, size_pt=14, bold=True, color=COLOR_WHITE)

        # 라벨 박스
        box = slide.shapes.add_shape(
            1,  # RECTANGLE
            x_pos, y_pos,
            step_width, Inches(1.8),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        box.line.color.rgb = COLOR_ACCENT
        box.line.width = Pt(1.5)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)

        p_label = tf.paragraphs[0]
        p_label.alignment = PP_ALIGN.CENTER
        run_l = p_label.add_run()
        run_l.text = label
        _set_font(run_l, size_pt=12, bold=True, color=COLOR_HEADER_BG)

        if desc:
            p_desc = tf.add_paragraph()
            p_desc.alignment = PP_ALIGN.CENTER
            p_desc.space_before = Pt(6)
            run_d = p_desc.add_run()
            run_d.text = desc
            _set_font(run_d, size_pt=9, color=COLOR_DARK_GRAY)

        # 화살표 (마지막 제외)
        if i < num_steps - 1:
            arrow_x = x_pos + step_width
            arrow_y = y_pos + Inches(0.9)
            arrow = slide.shapes.add_shape(
                1,  # RECTANGLE as arrow indicator
                arrow_x + Inches(0.05), arrow_y,
                gap - Inches(0.1), Inches(0.05),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_ACCENT
            arrow.line.fill.background()


def _build_timeline_slide(slide, slide_data: dict):
    """타임라인형 슬라이드: 시간순 이벤트를 수평선 위에 표현."""
    content = slide_data.get("content", {})
    events = content if isinstance(content, list) else content.get("events", [])
    if not events:
        events = [{"date": "TBD", "title": "이벤트", "description": ""}]

    body_top = HEADER_HEIGHT + Inches(0.5)
    num_events = len(events)
    line_y = body_top + Inches(2.0)

    # 수평 타임라인 선
    line = slide.shapes.add_shape(
        1,  # RECTANGLE
        Inches(1.0), line_y,
        SLIDE_WIDTH - Inches(2.0), Inches(0.06),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_HEADER_BG
    line.line.fill.background()

    available_width = SLIDE_WIDTH - Inches(2.0)
    item_width = available_width / max(num_events, 1)

    for i, event in enumerate(events):
        if isinstance(event, str):
            event = {"date": "", "title": event, "description": ""}
        date_text = event.get("date", "")
        title_text = event.get("title", "")
        desc_text = event.get("description", "")

        x_center = Inches(1.0) + item_width * i + item_width / 2

        # 원형 마커
        marker = slide.shapes.add_shape(
            9,  # OVAL
            x_center - Inches(0.15), line_y - Inches(0.12),
            Inches(0.3), Inches(0.3),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = COLOR_ACCENT
        marker.line.fill.background()

        # 날짜 (위)
        date_box = slide.shapes.add_textbox(
            x_center - Inches(0.8), line_y - Inches(0.8),
            Inches(1.6), Inches(0.5),
        )
        tf = date_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = date_text
        _set_font(run, size_pt=10, bold=True, color=COLOR_ACCENT)

        # 제목 + 설명 (아래)
        content_box = slide.shapes.add_textbox(
            x_center - Inches(0.9), line_y + Inches(0.35),
            Inches(1.8), Inches(1.5),
        )
        tf2 = content_box.text_frame
        tf2.word_wrap = True
        p_title = tf2.paragraphs[0]
        p_title.alignment = PP_ALIGN.CENTER
        run_t = p_title.add_run()
        run_t.text = title_text
        _set_font(run_t, size_pt=11, bold=True, color=COLOR_HEADER_BG)

        if desc_text:
            p_desc = tf2.add_paragraph()
            p_desc.alignment = PP_ALIGN.CENTER
            p_desc.space_before = Pt(4)
            run_d = p_desc.add_run()
            run_d.text = desc_text
            _set_font(run_d, size_pt=9, color=COLOR_DARK_GRAY)


def _build_hybrid_slide(slide, slide_data: dict):
    """하이브리드형 슬라이드: 왼쪽 텍스트 + 오른쪽 차트."""
    content = slide_data.get("content", {})

    # 텍스트 부분 (왼쪽 절반)
    text_items = content.get("text", content.get("bullets", []))
    if not text_items:
        text_items = slide_data.get("subPoints", ["(내용 없음)"])

    body_top = HEADER_HEIGHT + Inches(0.4)
    half_width = (SLIDE_WIDTH - Inches(1.6)) / 2
    body_height = SLIDE_HEIGHT - HEADER_HEIGHT - FOOTER_HEIGHT - Inches(0.8)

    # 왼쪽: 텍스트
    txBox = slide.shapes.add_textbox(Inches(0.8), body_top, half_width - Inches(0.2), body_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(text_items):
        if isinstance(item, dict):
            item = item.get("text", str(item))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"  {item}"
        _set_font(run, size_pt=12, color=COLOR_DARK_GRAY)

    # 오른쪽: 차트 (간단 막대)
    chart_data = content.get("chart", {})
    labels = chart_data.get("labels", ["A", "B", "C"])
    values = chart_data.get("values", [30, 50, 20])
    right_x = Inches(0.8) + half_width + Inches(0.2)
    max_val = max(values) if values else 1
    bar_height = min(Inches(0.4), body_height / max(len(labels), 1) - Inches(0.1))
    max_bar_width = half_width - Inches(2.0)

    for i, (label, value) in enumerate(zip(labels, values)):
        y_pos = body_top + Inches(0.1) + (bar_height + Inches(0.15)) * i

        label_box = slide.shapes.add_textbox(right_x, y_pos, Inches(1.5), bar_height)
        tf2 = label_box.text_frame
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        run2 = p2.add_run()
        run2.text = str(label)
        _set_font(run2, size_pt=9, color=COLOR_DARK_GRAY)

        bar_width = max(Inches(0.2), int(max_bar_width * (value / max_val)))
        bar = slide.shapes.add_shape(
            1, right_x + Inches(1.6), y_pos, bar_width, bar_height,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_ACCENT
        bar.line.fill.background()

        val_box = slide.shapes.add_textbox(
            right_x + Inches(1.6) + bar_width + Inches(0.05), y_pos,
            Inches(0.6), bar_height,
        )
        tf3 = val_box.text_frame
        tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        run3.text = str(value)
        _set_font(run3, size_pt=9, bold=True, color=COLOR_HEADER_BG)


# 레이아웃별 빌더 매핑
LAYOUT_BUILDERS = {
    "텍스트형": _build_text_slide,
    "표형": _build_table_slide,
    "차트형": _build_chart_slide,
    "다이어그램형": _build_diagram_slide,
    "타임라인형": _build_timeline_slide,
    "하이브리드형": _build_hybrid_slide,
}


def step7_generate_pptx(session: dict, output_path: str = None) -> dict:
    """PPTX 파일을 생성한다."""
    # 콘텐츠 생성
    logger.info("슬라이드 콘텐츠를 생성 중...")
    slides_content = _generate_slide_content(session)
    session["slides"] = slides_content

    # 출력 경로 결정
    if not output_path:
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        output_path = str(OUTPUT_DIR / f"presentation_{session['id']}.pptx")
    else:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)

    session["outputPath"] = output_path

    # PPTX 생성
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 빈 레이아웃 사용
    blank_layout = prs.slide_layouts[6]  # blank layout

    for i, slide_data in enumerate(slides_content):
        slide = prs.slides.add_slide(blank_layout)

        title = slide_data.get("title", slide_data.get("headMessage", f"슬라이드 {i+1}"))
        layout_type = slide_data.get("layout", "텍스트형")

        # 헤더
        _add_header(slide, title)

        # 본문 (레이아웃별)
        builder = LAYOUT_BUILDERS.get(layout_type, _build_text_slide)
        try:
            builder(slide, slide_data)
        except Exception as e:
            logger.warning(f"슬라이드 {i+1} ({layout_type}) 빌드 중 오류, 텍스트형으로 대체: {e}")
            _build_text_slide(slide, slide_data)

        # 푸터
        _add_footer(slide, i + 1)

    prs.save(output_path)
    logger.info(f"PPTX 저장 완료: {output_path}")

    session["step"] = 7
    save_session(session)

    return make_response(
        step=7,
        status="done",
        message=f"PPTX를 생성했습니다: {output_path}",
        data={
            "outputPath": output_path,
            "totalSlides": len(slides_content),
            "slides": [
                {"slideNumber": s.get("slideNumber", i+1), "title": s.get("title", ""), "layout": s.get("layout", "")}
                for i, s in enumerate(slides_content)
            ],
        },
        next_step=8,
        needs_input=False,
        prompt="PPTX 생성이 완료되었습니다.",
    )


# ─── Step 8: 최종화 ───────────────────────────────────────────────────────────

def step8_finalize(session: dict) -> dict:
    """세션을 최종화하고 결과를 반환한다."""
    session["step"] = 8
    save_session(session)

    return make_response(
        step=8,
        status="done",
        message="장표 제작이 완료되었습니다.",
        data={
            "outputPath": session["outputPath"],
            "sessionId": session["id"],
            "totalSlides": len(session.get("slides", [])),
        },
        next_step=-1,
        needs_input=False,
        prompt=f"최종 파일: {session['outputPath']}",
    )


# ─── 단계 실행 디스패처 ───────────────────────────────────────────────────────

STEP_HANDLERS = {
    1: lambda s, inp, out: step1_input_analysis(s, inp),
    2: lambda s, inp, out: step2_clarification(s, inp),
    3: lambda s, inp, out: step3_head_messages(s),
    4: lambda s, inp, out: step4_user_confirmation(s, inp),
    5: lambda s, inp, out: step5_layout_selection(s),
    6: lambda s, inp, out: step6_layout_confirmation(s, inp),
    7: lambda s, inp, out: step7_generate_pptx(s, out),
    8: lambda s, inp, out: step8_finalize(s),
}


def run_single_shot(user_input: str, output_path: str) -> dict:
    """Single-shot 모드: 확인 단계(2,4,6)를 건너뛰고 자동 생성한다."""
    session = new_session()
    results = []

    # Step 1: 입력 분석
    logger.info("Step 1: 입력 분석...")
    r1 = step1_input_analysis(session, user_input)
    results.append(r1)

    # Step 2: 자동 명확화 (기본 답변으로)
    logger.info("Step 2: 자동 명확화...")
    auto_answer = "경영진/실무진 대상, 전략 브리핑 수준, 초안 단계"
    r2 = step2_clarification(session, auto_answer)
    results.append(r2)

    # Step 3: 헤드메시지 생성
    logger.info("Step 3: 헤드메시지 생성...")
    r3 = step3_head_messages(session)
    results.append(r3)

    # Step 4: 자동 확정
    session["step"] = 4

    # Step 5: 레이아웃 선택
    logger.info("Step 5: 레이아웃 선택...")
    r5 = step5_layout_selection(session)
    results.append(r5)

    # Step 6: 자동 확정
    session["step"] = 6

    # Step 7: PPTX 생성
    logger.info("Step 7: PPTX 생성...")
    r7 = step7_generate_pptx(session, output_path)
    results.append(r7)

    # Step 8: 최종화
    logger.info("Step 8: 최종화...")
    r8 = step8_finalize(session)
    results.append(r8)

    return r8


def run_session_step(session_id: str, step: int, user_input: str = "",
                     output_path: str = None) -> dict:
    """Session 모드: 지정된 단계 하나를 실행한다."""
    session = load_session(session_id)

    handler = STEP_HANDLERS.get(step)
    if not handler:
        return make_response(
            step=step,
            status="error",
            message=f"알 수 없는 단계입니다: {step}",
            data={},
            next_step=step,
            needs_input=True,
            prompt="올바른 단계 번호(1-8)를 지정하세요.",
        )

    return handler(session, user_input, output_path)


# ─── CLI ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="장표 제작 에이전트 - 컨텍스트로부터 PPTX 자동 생성",
    )
    parser.add_argument("--input", "-i", type=str, default="",
                        help="입력 텍스트 또는 파일 경로")
    parser.add_argument("--output", "-o", type=str, default="",
                        help="출력 PPTX 파일 경로")
    parser.add_argument("--session", "-s", type=str, default="",
                        help="세션 ID (세션 모드)")
    parser.add_argument("--step", type=int, default=0,
                        help="실행할 단계 번호 (1-8, 세션 모드)")

    args = parser.parse_args()

    try:
        if args.session:
            # Session mode
            if args.step < 1 or args.step > 8:
                print(json.dumps({
                    "step": 0,
                    "status": "error",
                    "message": "단계 번호는 1-8 사이여야 합니다.",
                    "data": {},
                    "nextStep": 1,
                    "needsInput": True,
                    "prompt": "--step 1 부터 시작하세요.",
                }, ensure_ascii=False, indent=2))
                sys.exit(1)

            result = run_session_step(
                session_id=args.session,
                step=args.step,
                user_input=args.input,
                output_path=args.output or None,
            )
        else:
            # Single-shot mode
            if not args.input:
                print(json.dumps({
                    "step": 0,
                    "status": "error",
                    "message": "--input 또는 --session 옵션이 필요합니다.",
                    "data": {},
                    "nextStep": 1,
                    "needsInput": True,
                    "prompt": "사용법: python create.py --input '컨텍스트' --output output/result.pptx",
                }, ensure_ascii=False, indent=2))
                sys.exit(1)

            output_path = args.output or str(OUTPUT_DIR / "result.pptx")
            result = run_single_shot(args.input, output_path)

        print(json.dumps(result, ensure_ascii=False, indent=2))

    except json.JSONDecodeError as e:
        print(json.dumps({
            "step": -1,
            "status": "error",
            "message": f"Claude API 응답 파싱 실패: {e}",
            "data": {},
            "nextStep": -1,
            "needsInput": False,
            "prompt": "",
        }, ensure_ascii=False, indent=2))
        sys.exit(1)
    except ValueError as e:
        print(json.dumps({
            "step": -1,
            "status": "error",
            "message": str(e),
            "data": {},
            "nextStep": -1,
            "needsInput": False,
            "prompt": "",
        }, ensure_ascii=False, indent=2))
        sys.exit(1)
    except Exception as e:
        logger.exception("예상치 못한 오류 발생")
        print(json.dumps({
            "step": -1,
            "status": "error",
            "message": f"오류: {e}",
            "data": {},
            "nextStep": -1,
            "needsInput": False,
            "prompt": "",
        }, ensure_ascii=False, indent=2))
        sys.exit(1)


if __name__ == "__main__":
    main()
