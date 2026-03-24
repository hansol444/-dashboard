"""
create_test_ppt.py — 후처리 테스트용 한국어 PPT 생성 스크립트
실행: python create_test_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

W = Inches(13.33)
H = Inches(7.5)


def add_title(slide, text, font_size=24):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    tf.text = text
    tf.paragraphs[0].runs[0].font.size = Pt(font_size)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)


def add_textbox(slide, text, left, top, width, height, font_size=12, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    return txBox


def add_table(slide, data, left, top, width, height, header_font=11, body_font=10):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for r_idx, row in enumerate(data):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = cell_text
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(header_font if r_idx == 0 else body_font)
            if r_idx == 0:
                run.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def create_test_ppt(output_path: str):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    blank_layout = prs.slide_layouts[6]  # 빈 레이아웃

    # ── 슬라이드 1: 표지 ──────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_layout)
    add_textbox(slide1, "2024년 채용 현황 및 전략", Inches(1), Inches(2.5), Inches(11), Inches(1.5), font_size=32, bold=True)
    add_textbox(slide1, "인사팀 | 2024년 Q1 리뷰", Inches(1), Inches(4.2), Inches(8), Inches(0.7), font_size=16)
    add_textbox(slide1, "※ 본 자료는 대외비입니다", Inches(1), Inches(6.5), Inches(8), Inches(0.5), font_size=9)

    # ── 슬라이드 2: 채용 현황 요약 ────────────────────────────────────
    slide2 = prs.slides.add_slide(blank_layout)
    add_title(slide2, "채용 현황 요약 — 2024 Q1")
    add_textbox(slide2, "핵심 성과 지표 (KPI)", Inches(0.5), Inches(1.4), Inches(5), Inches(0.5), font_size=14, bold=True)

    bullets = (
        "• 총 지원자 수: 1,240명 (전년 대비 +18%)\n"
        "• 면접 통과율: 32%\n"
        "• 최종 합격률: 8.4%\n"
        "• 채용 소요 시간: 평균 23일\n"
        "• 채용 비용: 1인당 KRW 100M"
    )
    add_textbox(slide2, bullets, Inches(0.5), Inches(2.0), Inches(5.5), Inches(3), font_size=12)

    add_textbox(slide2, "예산 집행 현황", Inches(7), Inches(1.4), Inches(5), Inches(0.5), font_size=14, bold=True)
    budget_text = (
        "• 총 채용 예산: 14억 원\n"
        "• 1분기 집행액: 3.2억 원\n"
        "• 잔여 예산: 10.8억 원\n"
        "• ROI 목표: 전년 대비 +15%"
    )
    add_textbox(slide2, budget_text, Inches(7), Inches(2.0), Inches(5.5), Inches(3), font_size=12)

    add_textbox(slide2, "* 예산 단위: 억 원 / 채용 목표 달성률 기준으로 분기별 재검토 예정", Inches(0.5), Inches(6.5), Inches(12), Inches(0.5), font_size=9)

    # ── 슬라이드 3: 월별 채용 실적 (표) ──────────────────────────────
    slide3 = prs.slides.add_slide(blank_layout)
    add_title(slide3, "월별 채용 실적 — 1월~3월")

    table_data = [
        ["구분", "1월", "2월", "3월", "합계"],
        ["지원자 수", "380명", "420명", "440명", "1,240명"],
        ["서류 합격", "120명", "138명", "148명", "406명"],
        ["면접 진행", "85명", "98명", "104명", "287명"],
        ["최종 합격", "28명", "34명", "42명", "104명"],
        ["합격률", "7.4%", "8.1%", "9.5%", "8.4%"],
    ]
    add_table(
        slide3, table_data,
        Inches(0.5), Inches(1.5), Inches(12), Inches(4.5),
        header_font=11, body_font=10
    )
    add_textbox(slide3, "※ 1월 실적은 신년 채용 캠페인 영향으로 지원자 수 감소", Inches(0.5), Inches(6.3), Inches(10), Inches(0.5), font_size=9)

    # ── 슬라이드 4: 채용 채널별 분석 ─────────────────────────────────
    slide4 = prs.slides.add_slide(blank_layout)
    add_title(slide4, "채용 채널별 분석")

    channel_data = [
        ["채널", "지원자 수", "합격자 수", "합격률", "채용 비용"],
        ["LinkedIn", "520명", "44명", "8.5%", "KRW 45M"],
        ["Seek", "310명", "28명", "9.0%", "KRW 30M"],
        ["사내 추천", "180명", "22명", "12.2%", "KRW 8M"],
        ["Indeed", "150명", "6명", "4.0%", "KRW 12M"],
        ["기타", "80명", "4명", "5.0%", "KRW 5M"],
    ]
    add_table(
        slide4, channel_data,
        Inches(0.5), Inches(1.5), Inches(12), Inches(4.2),
        header_font=11, body_font=10
    )

    insight = (
        "핵심 인사이트:\n"
        "• 사내 추천 채널의 합격률이 가장 높음 (12.2%)\n"
        "• LinkedIn은 지원자 볼륨은 크나 비용 효율 개선 필요\n"
        "• APAC 전략 차원에서 ANZ 채널 다각화 검토 중\n"
        "• HRBP와 협력해 OKR 기반 채용 목표 재설정 예정"
    )
    add_textbox(slide4, insight, Inches(0.5), Inches(6.0), Inches(12), Inches(1.2), font_size=10)

    # ── 슬라이드 5: Q2 전략 및 실행 계획 ─────────────────────────────
    slide5 = prs.slides.add_slide(blank_layout)
    add_title(slide5, "Q2 채용 전략 및 실행 계획")

    strategy = (
        "전략 방향\n"
        "• 인재풀 확대: 타겟 직군 지원자 20% 증가 목표\n"
        "• 고용주 브랜딩 강화: SNS 채용 콘텐츠 월 8회 발행\n"
        "• 채용 프로세스 개선: 채용 소요 시간 23일 → 18일 단축\n"
        "• 온보딩 만족도 목표: 85점 이상 (현재 79점)"
    )
    add_textbox(slide5, strategy, Inches(0.5), Inches(1.5), Inches(6), Inches(4), font_size=12)

    action_plan = (
        "분기별 실행 계획\n"
        "• 4월: 직무기술서 전면 개편, 채용 공고 최적화\n"
        "• 5월: 면접관 교육 프로그램 도입 (조직문화 중심)\n"
        "• 6월: 상반기 채용 결과 리뷰, 하반기 인력 계획 수립\n"
        "• 예산: 분기 총 3.5억 원 (전분기 대비 +9.4%)"
    )
    add_textbox(slide5, action_plan, Inches(7), Inches(1.5), Inches(5.8), Inches(4), font_size=12)

    add_textbox(slide5, "담당: 인사팀 채용파트 | 검토: CHRO | 승인: CEO", Inches(0.5), Inches(6.4), Inches(10), Inches(0.5), font_size=9)

    os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
    prs.save(output_path)
    print(f"[OK] 테스트 PPT 생성 완료: {output_path}")
    print(f"     슬라이드 수: 5장")
    print(f"     후처리 테스트 요소: 14억(억 단위), 1월~3월(월 표기), KRW 100M(통화 순서), KPI/APAC/ANZ(보존)")


if __name__ == "__main__":
    create_test_ppt("input/테스트_HR_채용.pptx")
