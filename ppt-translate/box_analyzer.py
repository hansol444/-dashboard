"""
box_analyzer.py — Q1: 텍스트박스 크기 분석 및 영문 최대 글자수 계산
"""

import math
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 1 Emu = 1/914400 inch = 1/914400 * 72 pt
EMU_PER_PT = 914400 / 72
DEFAULT_FONT_SIZE_PT = 12
# 영문 글자 평균 너비 계수 (폰트 크기 대비)
CHAR_WIDTH_RATIO = 0.6


def emu_to_pt(emu: int) -> float:
    return emu / EMU_PER_PT


def calc_max_chars(box_width_emu: int, font_size_pt: float) -> int:
    """텍스트박스 너비와 폰트 크기로 들어갈 수 있는 영문 최대 글자수 계산."""
    box_width_pt = emu_to_pt(box_width_emu)
    char_width_pt = font_size_pt * CHAR_WIDTH_RATIO
    if char_width_pt <= 0:
        return 999
    return math.floor(box_width_pt / char_width_pt)


def get_font_size_from_run(run) -> float:
    """런에서 폰트 크기를 가져오고, 없으면 None 반환."""
    if run.font and run.font.size:
        return run.font.size / 12700  # EMU → pt (1pt = 12700 EMU in pptx)
    return None


def get_font_size_from_paragraph(para) -> float:
    """단락에서 폰트 크기를 추출 (런 순서대로 첫 번째 유효값 사용)."""
    for run in para.runs:
        size = get_font_size_from_run(run)
        if size:
            return size
    return None


def get_dominant_font_size(text_frame) -> float:
    """텍스트 프레임 전체에서 가장 많이 쓰인 폰트 크기 반환."""
    sizes = []
    for para in text_frame.paragraphs:
        size = get_font_size_from_paragraph(para)
        if size:
            sizes.append(size)
    if not sizes:
        return DEFAULT_FONT_SIZE_PT
    # 가장 빈번한 크기 반환
    return max(set(sizes), key=sizes.count)


def _is_smartart(shape) -> bool:
    """shape이 SmartArt(dgm) 도형인지 확인."""
    DGM_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
    try:
        for elem in shape._element.iter():
            if elem.get("uri") == DGM_URI:
                return True
    except Exception:
        pass
    return False


def _extract_smartart_blocks(shape, slide_idx: int, slide_part) -> list[dict]:
    """SmartArt 다이어그램 데이터에서 텍스트 블록 추출."""
    DGM_NS = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
    A_NS   = "http://schemas.openxmlformats.org/drawingml/2006/main"
    R_NS   = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    blocks = []
    try:
        from lxml import etree
        elem = shape._element
        relIds_elem = elem.find(f'.//{{{DGM_NS}}}relIds')
        if relIds_elem is None:
            return blocks
        dm_rId = relIds_elem.get(f'{{{R_NS}}}dm')
        if not dm_rId:
            return blocks
        dgm_part = slide_part.related_parts[dm_rId]
        dgm_root = etree.fromstring(dgm_part.blob)
        node_idx = 0
        for t_node in dgm_root.iter(f'{{{A_NS}}}t'):
            if t_node.text and t_node.text.strip():
                blocks.append({
                    "slide_idx": slide_idx,
                    "shape_id": shape.shape_id,
                    "shape_name": f"{shape.name} [스마트아트 {node_idx}]",
                    "shape_type": "smartart",
                    "text": t_node.text.strip(),
                    "smartart_node_index": node_idx,
                    "font_size_pt": DEFAULT_FONT_SIZE_PT,
                    "box_width_emu": shape.width,
                    "box_width_pt": round(emu_to_pt(shape.width), 2),
                    "max_chars": 999,
                })
                node_idx += 1
    except Exception:
        pass
    return blocks


def extract_text_blocks_from_shape(shape, slide_idx: int, slide_part=None) -> list[dict]:
    """단일 shape에서 텍스트 블록 목록을 추출."""
    blocks = []

    # 그룹 도형: 자식 도형으로 재귀
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            blocks.extend(extract_text_blocks_from_shape(child, slide_idx, slide_part))
        return blocks

    # 차트: 제목만 추출
    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        try:
            chart = shape.chart
            if chart.has_title and chart.chart_title.has_text_frame:
                text = chart.chart_title.text_frame.text.strip()
                if text:
                    blocks.append({
                        "slide_idx": slide_idx,
                        "shape_id": shape.shape_id,
                        "shape_name": shape.name + " [차트제목]",
                        "shape_type": "chart_title",
                        "text": text,
                        "font_size_pt": DEFAULT_FONT_SIZE_PT,
                        "box_width_emu": shape.width,
                        "box_width_pt": round(emu_to_pt(shape.width), 2),
                        "max_chars": 999,
                    })
        except Exception:
            pass
        return blocks

    # SmartArt
    if _is_smartart(shape) and slide_part is not None:
        return _extract_smartart_blocks(shape, slide_idx, slide_part)

    # 일반 텍스트박스 / 제목 / 콘텐츠 placeholder
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if not text:
            return blocks

        font_size = get_dominant_font_size(shape.text_frame)
        max_chars = calc_max_chars(shape.width, font_size)

        # 다중 줄 텍스트: 줄 수만큼 max_chars를 확대
        line_count = max(text.count("\n") + 1, 1)
        total_max_chars = max_chars * line_count

        blocks.append({
            "slide_idx": slide_idx,
            "shape_id": shape.shape_id,
            "shape_name": shape.name,
            "shape_type": "textbox",
            "text": text,
            "font_size_pt": font_size,
            "box_width_emu": shape.width,
            "box_width_pt": round(emu_to_pt(shape.width), 2),
            "max_chars": total_max_chars,
        })

    # 표(Table)
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table = shape.table
        col_count = len(table.columns)
        # 표 셀 너비: 표 전체 너비를 열 수로 균등 분배 (근사값)
        cell_width_emu = shape.width // col_count if col_count > 0 else shape.width

        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                text = cell.text_frame.text.strip()
                if not text:
                    continue
                font_size = get_dominant_font_size(cell.text_frame)
                # 병합 셀은 실제 너비가 더 넓지만 근사값 사용
                max_chars = calc_max_chars(cell_width_emu, font_size)

                blocks.append({
                    "slide_idx": slide_idx,
                    "shape_id": shape.shape_id,
                    "shape_name": f"{shape.name} [행{row_idx} 열{col_idx}]",
                    "shape_type": "table_cell",
                    "row": row_idx,
                    "col": col_idx,
                    "text": text,
                    "font_size_pt": font_size,
                    "box_width_emu": cell_width_emu,
                    "box_width_pt": round(emu_to_pt(cell_width_emu), 2),
                    "max_chars": max_chars,
                })

    return blocks


def analyze_pptx(pptx_path: str) -> dict:
    """
    PPT 파일을 분석해 슬라이드별 텍스트 블록 정보를 반환.

    반환값:
        {
            slide_idx (int): [
                {
                    slide_idx, shape_id, shape_name, shape_type,
                    text, font_size_pt, box_width_emu, box_width_pt, max_chars,
                    [row, col]  # 표 셀인 경우만
                },
                ...
            ]
        }
    """
    prs = Presentation(pptx_path)
    result = {}

    for slide_idx, slide in enumerate(prs.slides):
        blocks = []
        for shape in slide.shapes:
            blocks.extend(extract_text_blocks_from_shape(shape, slide_idx, slide.part))

        # 슬라이드 노트
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                # 노트는 너비 제약을 슬라이드 너비로 간주
                font_size = DEFAULT_FONT_SIZE_PT
                max_chars = calc_max_chars(prs.slide_width, font_size)
                blocks.append({
                    "slide_idx": slide_idx,
                    "shape_id": -1,
                    "shape_name": "notes",
                    "shape_type": "notes",
                    "text": notes_text,
                    "font_size_pt": font_size,
                    "box_width_emu": prs.slide_width,
                    "box_width_pt": round(emu_to_pt(prs.slide_width), 2),
                    "max_chars": max_chars,
                })

        if blocks:
            result[slide_idx] = blocks

    return result


if __name__ == "__main__":
    import sys
    import json

    path = sys.argv[1] if len(sys.argv) > 1 else "input/테스트_HR_채용.pptx"
    data = analyze_pptx(path)
    for slide_idx, blocks in data.items():
        print(f"\n=== 슬라이드 {slide_idx + 1} ===")
        for b in blocks:
            print(f"  [{b['shape_name']}] {b['font_size_pt']}pt | max {b['max_chars']}자 | {b['text'][:40]}")
