"""
translate.py — PPT 장표 번역 자동화 메인 스크립트

사용법:
  python translate.py input/파일.pptx --to en
  python translate.py input/파일.pptx --to ko
  python translate.py --batch --to en
  python translate.py input/파일.pptx --quality precise
  python translate.py input/파일.pptx --no-postprocess
  python translate.py input/파일.pptx --no-report
"""

import sys, os

def _safe_print(*args, **kwargs):
    """UnicodeEncodeError 없이 안전하게 출력."""
    try:
        print(*args, **kwargs)
    except UnicodeEncodeError:
        safe = " ".join(str(a) for a in args).encode("ascii", errors="replace").decode("ascii")
        try:
            print(safe, **kwargs)
        except Exception:
            pass
    except Exception:
        pass

import argparse
import json
import time
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

import anthropic
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import config
from box_analyzer import analyze_pptx
from post_processor import PostProcessor


# ─── API 클라이언트 초기화 ─────────────────────────────────────────────────────

def get_api_client() -> anthropic.Anthropic:
    api_key = config.ANTHROPIC_API_KEY
    if not api_key:
        raise ValueError(
            "ANTHROPIC_API_KEY가 설정되지 않았습니다.\n"
            "환경변수: set ANTHROPIC_API_KEY=your_key_here\n"
            "또는 config.py에 직접 입력하세요."
        )
    return anthropic.Anthropic(api_key=api_key)


# ─── 용어집 로드 ───────────────────────────────────────────────────────────────

def load_terminology(direction: str) -> tuple[dict, list]:
    """terminology.json 로드. direction: 'ko_to_en' or 'en_to_ko'"""
    try:
        with open(config.TERMINOLOGY_PATH, encoding="utf-8") as f:
            data = json.load(f)
        terms = data.get(direction, {})
        preserve = data.get("preserve", [])
        return terms, preserve
    except FileNotFoundError:
        return {}, []


def build_glossary_prompt(terms: dict, preserve: list, slide_text: str = "") -> str:
    """
    슬라이드 텍스트에 등장하는 용어만 필터링해서 프롬프트 생성.
    slide_text가 없으면 전체 용어집 사용.
    """
    lines = []

    if terms:
        if slide_text:
            # 정확한 단어 경계 매칭: "채용"이 "채용 시장"에 매칭되되,
            # "개인"이 "개인화"에 오매칭되지 않도록 함
            import re
            filtered = {}
            for ko, en in terms.items():
                # 한글은 \b가 안 먹으므로, 용어 앞뒤가 다른 한글이 아닌지 확인
                pattern = re.compile(r'(?<![가-힣])' + re.escape(ko) + r'(?![가-힣])')
                if pattern.search(slide_text):
                    filtered[ko] = en
        else:
            filtered = terms

        if filtered:
            lines.append("## Glossary (use these translations exactly)")
            for ko, en in filtered.items():
                lines.append(f"- {ko} → {en}")

    if preserve:
        # preserve는 전체 유지 (짧은 약어라 토큰 부담 적음)
        lines.append("\n## Do NOT translate these terms (preserve as-is)")
        lines.append(", ".join(preserve))

    return "\n".join(lines)


# ─── 시스템 프롬프트 로드 ──────────────────────────────────────────────────────

def load_base_prompt(direction: str) -> str:
    """용어집 없이 기본 시스템 프롬프트만 로드."""
    try:
        with open(config.SYSTEM_PROMPT_PATH, encoding="utf-8") as f:
            base_prompt = f.read()
    except FileNotFoundError:
        base_prompt = "You are a professional translator."

    if direction == "en_to_ko":
        base_prompt = base_prompt.replace(
            "Translate Korean to English", "Translate English to Korean"
        )
    return base_prompt


def build_slide_prompt(base_prompt: str, terms: dict, preserve: list, slide_text: str) -> str:
    """슬라이드 텍스트 기준으로 필터링된 최종 시스템 프롬프트 생성."""
    glossary = build_glossary_prompt(terms, preserve, slide_text)
    return f"{base_prompt}\n\n{glossary}" if glossary else base_prompt


# ─── Claude API 번역 호출 ──────────────────────────────────────────────────────

# 병렬 호출 시 rate limit 동시 재시도 방지용 락
_rate_limit_lock = threading.Lock()


def call_api(client: anthropic.Anthropic, system_prompt: str, user_message: str) -> str:
    """API 호출 with 재시도. rate limit 시 스레드 간 직렬화.
    system_prompt에 cache_control을 걸어 동일 프롬프트 반복 호출 시 비용/속도 절감.
    """
    for attempt in range(config.API_MAX_RETRIES):
        try:
            response = client.messages.create(
                model=config.MODEL,
                max_tokens=4096,
                system=[{
                    "type": "text",
                    "text": system_prompt,
                    "cache_control": {"type": "ephemeral"},
                }],
                messages=[{"role": "user", "content": user_message}],
            )
            return response.content[0].text.strip()
        except anthropic.APIConnectionError:
            if attempt < config.API_MAX_RETRIES - 1:
                _safe_print(f"  [!] 네트워크 오류, {config.API_RETRY_DELAY}초 후 재시도... ({attempt+1}/{config.API_MAX_RETRIES})")
                time.sleep(config.API_RETRY_DELAY)
            else:
                raise
        except anthropic.RateLimitError:
            with _rate_limit_lock:
                wait = config.API_RETRY_DELAY * (attempt + 1) * 2
                _safe_print(f"  [!] Rate limit, {wait}초 대기...")
                time.sleep(wait)
    return ""


def translate_text(
    client: anthropic.Anthropic,
    system_prompt: str,
    slide_context: str,
    text: str,
    max_chars: int,
    allow_abbreviation: bool = False,
) -> str:
    """단일 텍스트 번역 요청."""
    constraint = f"영문 {max_chars}자 이내로 번역하세요." if max_chars < 999 else ""
    if allow_abbreviation:
        constraint += " 공간이 부족하므로 관사 생략, 단어 축약이 허용됩니다."

    user_msg = f"""[슬라이드 전체 맥락]
{slide_context}

[번역할 텍스트]
{text}

{constraint}
번역문만 출력하세요. 설명 없이."""

    return call_api(client, system_prompt, user_msg)


def translate_slide_batch(
    client: anthropic.Anthropic,
    system_prompt: str,
    slide_context: str,
    blocks: list[dict],
) -> list[str]:
    """슬라이드 전체 텍스트를 한 번에 번역 (fast 모드)."""
    numbered = "\n".join(
        f"[{i+1}] (max {b['max_chars']}자) {b['text']}"
        for i, b in enumerate(blocks)
    )
    user_msg = f"""[슬라이드 전체 맥락]
{slide_context}

[번역할 텍스트 목록 — 번호 순서 그대로 같은 번호로 반환하세요]
{numbered}

각 항목을 번호 포함해서 아래 형식으로 출력:
[1] 번역문
[2] 번역문
...

설명 없이 번역문만."""

    raw = call_api(client, system_prompt, user_msg)

    # 파싱: [N] 패턴으로 분리
    import re
    results = {}
    for match in re.finditer(r'\[(\d+)\]\s*(.*?)(?=\[\d+\]|$)', raw, re.DOTALL):
        idx = int(match.group(1)) - 1
        text = match.group(2).strip()
        if text:   # 빈 번역 무시 → 원문 유지
            results[idx] = text

    return [results.get(i, blocks[i]["text"]) for i in range(len(blocks))]


# ─── 서식 보존 텍스트 삽입 ────────────────────────────────────────────────────

def apply_text_to_shape(shape, translated_text: str, font_size_pt: float = None):
    """shape의 텍스트를 번역문으로 교체. font_size_pt 지정 시 모든 run에 해당 크기 적용."""
    if not shape.has_text_frame:
        return
    if not translated_text.strip():   # 빈 번역 → 원문 유지
        return

    tf = shape.text_frame
    lines = translated_text.split("\n")
    para_count = len(tf.paragraphs)

    # 번역문 줄이 단락보다 많으면 마지막 단락에 나머지를 공백으로 이어붙임
    # (PPT run 안의 \n은 줄바꿈으로 렌더링되지 않으므로 공백 연결)
    if len(lines) > para_count and para_count > 0:
        lines = lines[:para_count - 1] + [" ".join(lines[para_count - 1:])]

    for para_idx, para in enumerate(tf.paragraphs):
        if para_idx >= len(lines):
            # 남은 단락 비우기 (단락 서식은 유지)
            for run in para.runs:
                run.text = ""
            continue

        line_text = lines[para_idx]

        if para.runs:
            para.runs[0].text = line_text
            if font_size_pt is not None:
                from pptx.util import Pt
                para.runs[0].font.size = Pt(font_size_pt)
            for run in para.runs[1:]:
                run.text = ""
        else:
            from pptx.oxml.ns import qn
            from lxml import etree
            r = etree.SubElement(para._p, qn('a:r'))
            rPr = etree.SubElement(r, qn('a:rPr'), attrib={'lang': 'en-AU'})
            t = etree.SubElement(r, qn('a:t'))
            t.text = line_text


def apply_text_to_table_cell(cell, translated_text: str, font_size_pt: float = None):
    """표 셀의 텍스트를 번역문으로 교체 (다단락 셀 포함)."""
    if not translated_text.strip():
        return
    tf = cell.text_frame
    lines = translated_text.split("\n")
    para_count = len(tf.paragraphs)

    # 번역문 줄이 단락보다 많으면 마지막 단락에 합침
    if len(lines) > para_count and para_count > 0:
        lines = lines[:para_count - 1] + [" ".join(lines[para_count - 1:])]

    for para_idx, para in enumerate(tf.paragraphs):
        if para_idx >= len(lines):
            # 남는 단락 비우기 (원문 한국어가 남지 않도록)
            for run in para.runs:
                run.text = ""
            continue
        line_text = lines[para_idx]
        if para.runs:
            para.runs[0].text = line_text
            if font_size_pt is not None:
                from pptx.util import Pt
                para.runs[0].font.size = Pt(font_size_pt)
            for run in para.runs[1:]:
                run.text = ""


# ─── 그룹/차트/SmartArt 헬퍼 ────────────────────────────────────────────────────

def _build_shape_map(shapes) -> dict:
    """slide.shapes를 재귀 순회해 {shape_id: shape} 맵 생성 (그룹 내부 포함)."""
    result = {}
    for shape in shapes:
        result[shape.shape_id] = shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            result.update(_build_shape_map(shape.shapes))
    return result


def apply_text_to_chart_title(shape, translated_text: str):
    """차트 제목 텍스트 교체."""
    if not translated_text.strip():
        return
    try:
        tf = shape.chart.chart_title.text_frame
        if tf.paragraphs and tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].text = translated_text
            for run in tf.paragraphs[0].runs[1:]:
                run.text = ""
    except Exception:
        pass


def _apply_smartart_node(shape, node_index: int, translated_text: str, slide_part):
    """SmartArt 다이어그램 데이터의 특정 텍스트 노드를 번역문으로 교체."""
    DGM_NS = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
    A_NS   = "http://schemas.openxmlformats.org/drawingml/2006/main"
    R_NS   = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    try:
        from lxml import etree
        elem = shape._element
        relIds_elem = elem.find(f'.//{{{DGM_NS}}}relIds')
        if relIds_elem is None:
            return
        dm_rId = relIds_elem.get(f'{{{R_NS}}}dm')
        if not dm_rId:
            return
        dgm_part = slide_part.related_parts[dm_rId]
        dgm_root = etree.fromstring(dgm_part.blob)
        t_nodes = [n for n in dgm_root.iter(f'{{{A_NS}}}t') if n.text and n.text.strip()]
        if node_index < len(t_nodes):
            t_nodes[node_index].text = translated_text
            dgm_part.blob = etree.tostring(dgm_root)
    except Exception as e:
        _safe_print(f"  [!] SmartArt 번역 적용 실패: {e}")


# ─── 번역 리포트 ───────────────────────────────────────────────────────────────

class TranslationReport:
    def __init__(self):
        self.entries = []
        self.overflow_count = 0
        self.postprocess_changes = []
        self.start_time = time.time()

    def add(self, slide_idx, original, translated, max_chars, overflow, pp_changes):
        status = "[!] 초과" if overflow else "OK"
        self.entries.append({
            "slide": slide_idx + 1,
            "original": original,
            "translated": translated,
            "max_chars": max_chars,
            "actual_chars": len(translated),
            "status": status,
        })
        if overflow:
            self.overflow_count += 1
        self.postprocess_changes.extend(pp_changes)

    def save(self, output_path: str):
        elapsed = round(time.time() - self.start_time, 1)
        total = len(self.entries)
        lines = [
            "=" * 70,
            "PPT 번역 리포트",
            "=" * 70,
            f"총 번역 블록: {total}  |  초과 항목: {self.overflow_count}  |  소요 시간: {elapsed}초",
            "",
            f"{'슬라이드':<6} {'원문':<30} {'번역문':<30} {'글자제약':<6} {'실제':<6} {'상태'}",
            "-" * 70,
        ]
        for e in self.entries:
            orig = e["original"][:28] + ".." if len(e["original"]) > 30 else e["original"]
            trans = e["translated"][:28] + ".." if len(e["translated"]) > 30 else e["translated"]
            lines.append(
                f"{e['slide']:<6} {orig:<30} {trans:<30} {e['max_chars']:<6} {e['actual_chars']:<6} {e['status']}"
            )

        if self.postprocess_changes:
            lines += ["", "─" * 70, "후처리 수정 목록", "─" * 70]
            lines += self.postprocess_changes

        lines += ["", "=" * 70]

        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))
        _safe_print(f"  [리포트] 저장: {output_path}")


# ─── 메인 번역 로직 ────────────────────────────────────────────────────────────

def translate_pptx(
    input_path: str,
    output_path: str,
    direction: str = "ko_to_en",
    quality: str = "fast",
    postprocess: bool = True,
    make_report: bool = True,
    enabled_rules: list = None,
    progress_callback=None,  # callable(event_type: str, data: dict)
):
    analysis = analyze_pptx(input_path)
    total_slides = len(analysis)

    if total_slides == 0:
        return

    client = get_api_client()
    terms, preserve = load_terminology(direction)
    base_prompt = load_base_prompt(direction)
    rules = enabled_rules if enabled_rules is not None else config.get_enabled_rules()
    pp = PostProcessor(rules) if postprocess else None
    report = TranslationReport() if make_report else None

    prs = Presentation(input_path)
    start = time.time()

    all_slides = list(analysis.keys())
    total_slides = len(all_slides)

    # ── Phase 1: 모든 슬라이드 API 호출을 병렬로 실행 ────────────────────────
    def translate_one_slide(args):
        slide_num, slide_idx = args
        blocks = analysis[slide_idx]
        slide_context = "\n".join(b["text"] for b in blocks)
        # 이 슬라이드 텍스트에 등장하는 용어만 필터링한 프롬프트
        slide_prompt = build_slide_prompt(base_prompt, terms, preserve, slide_context)
        if quality == "fast":
            translations = translate_slide_batch(client, slide_prompt, slide_context, blocks)
        else:
            translations = [
                translate_text(client, slide_prompt, slide_context, b["text"], b["max_chars"])
                for b in blocks
            ]
        return slide_num, slide_idx, blocks, translations, slide_context, slide_prompt

    slide_results = {}
    completed_count = 0
    max_workers = min(total_slides, config.MAX_PARALLEL_SLIDES)

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = {
            executor.submit(translate_one_slide, (slide_num, slide_idx)): slide_num
            for slide_num, slide_idx in enumerate(all_slides)
            if analysis[slide_idx]
        }
        for future in as_completed(futures):
            try:
                slide_num, slide_idx, blocks, translations, slide_context, slide_prompt = future.result()
            except Exception as e:
                # 개별 슬라이드 실패 시 건너뛰고 나머지 계속 진행
                completed_count += 1
                _safe_print(f"\n  [!] 슬라이드 번역 실패: {e}")
                continue
            slide_results[slide_idx] = (slide_num, blocks, translations, slide_context, slide_prompt)
            completed_count += 1
            elapsed = time.time() - start
            ratio = completed_count / total_slides
            eta = (elapsed / ratio - elapsed) if ratio > 0 else 0
            if progress_callback:
                progress_callback("progress", {
                    "slide": completed_count,
                    "total": total_slides,
                    "percent": int(ratio * 100),
                    "elapsed": int(elapsed),
                    "eta": int(eta),
                })
            else:
                bar = "#" * int(ratio * 10) + "." * (10 - int(ratio * 10))
                _safe_print(f"\r  {bar} {completed_count}/{total_slides} 번역 완료 | {int(elapsed)}초   ", end="", flush=True)

    # ── Phase 2a: 후처리 + 초과 항목 수집 ───────────────────────────────────
    final = {}      # {(slide_idx, block_idx): (original, translated, block, pp_changes, overflow)}
    retry_list = [] # [(slide_idx, block_idx, block, original, slide_prompt, slide_context)]

    for slide_idx in all_slides:
        if slide_idx not in slide_results:
            continue
        _, blocks, translations, slide_context, slide_prompt = slide_results[slide_idx]
        for bi, (block, translated) in enumerate(zip(blocks, translations)):
            original = block["text"]
            pp_changes = []
            if pp:
                translated, pp_changes = pp.process(translated, original)
            overflow = len(translated) > block["max_chars"]
            final[(slide_idx, bi)] = (original, translated, block, pp_changes, overflow)
            if overflow:
                retry_list.append((slide_idx, bi, block, original, slide_prompt, slide_context))

    # ── Phase 2b: 초과 항목 병렬 재번역 ─────────────────────────────────────
    total_retry = len(retry_list)
    retry_done  = 0

    def _do_retry(task):
        s_idx, b_idx, block, original, slide_prompt, slide_context = task
        retry = translate_text(
            client, slide_prompt, slide_context,
            original, block["max_chars"], allow_abbreviation=True,
        )
        changes = []
        if pp:
            retry, changes = pp.process(retry, original)
        return s_idx, b_idx, retry, changes, len(retry) > block["max_chars"]

    if retry_list:
        if progress_callback:
            progress_callback("status", {"message": f"초과 항목 재번역 중 (0/{total_retry})..."})
        with ThreadPoolExecutor(max_workers=min(total_retry, config.MAX_PARALLEL_SLIDES)) as ex:
            futs = {ex.submit(_do_retry, t): t for t in retry_list}
            for fut in as_completed(futs):
                try:
                    s_idx, b_idx, retry, changes, overflow = fut.result()
                    orig, _, blk, old_changes, _ = final[(s_idx, b_idx)]
                    final[(s_idx, b_idx)] = (orig, retry, blk, old_changes + changes, overflow)
                except Exception:
                    pass
                retry_done += 1
                if progress_callback:
                    progress_callback("status", {"message": f"초과 항목 재번역 중 ({retry_done}/{total_retry})..."})

    # ── Phase 2c: PPT에 적용 (단일 스레드) ───────────────────────────────────
    for slide_idx in all_slides:
        if slide_idx not in slide_results:
            continue
        _, blocks, _, _, _ = slide_results[slide_idx]
        slide     = prs.slides[slide_idx]
        shape_map = _build_shape_map(slide.shapes)  # 그룹 내부 포함

        for bi, block in enumerate(blocks):
            key = (slide_idx, bi)
            if key not in final:
                continue
            original, translated, block, pp_changes, overflow = final[key]

            if overflow and progress_callback:
                progress_callback("overflow", {
                    "slide": slide_idx + 1,
                    "text": original[:60],
                    "max_chars": block["max_chars"],
                    "actual_chars": len(translated),
                })

            shape = shape_map.get(block["shape_id"])
            shape_type = block["shape_type"]

            # overflow 시 글자 크기 축소 (빨간색 대신)
            font_size_pt = None
            if overflow and block.get("font_size_pt") and block.get("max_chars", 0) > 0:
                actual = max(len(translated), 1)
                ratio = block["max_chars"] / actual
                font_size_pt = max(round(block["font_size_pt"] * ratio, 1), 6.0)

            if shape is not None:
                if shape_type == "table_cell":
                    cell = shape.table.cell(block["row"], block["col"])
                    apply_text_to_table_cell(cell, translated, font_size_pt)
                elif shape_type == "chart_title":
                    apply_text_to_chart_title(shape, translated)
                elif shape_type == "smartart":
                    _apply_smartart_node(shape, block.get("smartart_node_index", 0), translated, slide.part)
                elif shape_type != "notes":
                    apply_text_to_shape(shape, translated, font_size_pt)
            elif shape_type == "notes" and slide.has_notes_slide:
                tf = slide.notes_slide.notes_text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = translated

            if report:
                report.add(slide_idx, original, translated, block["max_chars"], overflow, pp_changes)

    if not progress_callback:
        _safe_print()

    os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
    prs.save(output_path)
    if not progress_callback:
        _safe_print(f"  [완료] 저장: {output_path}")

    if report:
        report_path = output_path.replace(".pptx", "_번역리포트.txt")
        report.save(report_path)


# ─── CLI ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="PPT 장표 번역 자동화")
    parser.add_argument("input", nargs="?", help="번역할 .pptx 파일 경로")
    parser.add_argument("--to", choices=["en", "ko"], default=config.TARGET_LANGUAGE, help="번역 방향 (기본: en)")
    parser.add_argument("--batch", action="store_true", help=f"{config.INPUT_DIR}/ 폴더 내 전체 파일 번역")
    parser.add_argument("--quality", choices=["fast", "precise"], default=config.QUALITY_MODE)
    parser.add_argument("--no-postprocess", action="store_true", help="후처리 비활성화")
    parser.add_argument("--no-report", action="store_true", help="번역 리포트 생성 안 함")
    args = parser.parse_args()

    direction = "ko_to_en" if args.to == "en" else "en_to_ko"
    postprocess = not args.no_postprocess
    make_report = not args.no_report

    if args.batch:
        input_dir = Path(config.INPUT_DIR)
        files = list(input_dir.glob("*.pptx"))
        if not files:
            _safe_print(f"[오류] {config.INPUT_DIR}/ 폴더에 .pptx 파일이 없습니다.")
            sys.exit(1)
        _safe_print(f"[배치] 번역: {len(files)}개 파일")
        for f in files:
            suffix = "_EN" if args.to == "en" else "_KO"
            out = str(Path(config.OUTPUT_DIR) / (f.stem + suffix + ".pptx"))
            translate_pptx(str(f), out, direction, args.quality, postprocess, make_report)
    elif args.input:
        if not args.input.endswith(".pptx"):
            _safe_print(f"[오류] .pptx 파일만 지원됩니다: {args.input}")
            sys.exit(1)
        if not Path(args.input).exists():
            _safe_print(f"[오류] 파일을 찾을 수 없습니다: {args.input}")
            sys.exit(1)
        suffix = "_EN" if args.to == "en" else "_KO"
        stem = Path(args.input).stem
        out = str(Path(config.OUTPUT_DIR) / (stem + suffix + ".pptx"))
        translate_pptx(args.input, out, direction, args.quality, postprocess, make_report)
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
