# -*- coding: utf-8 -*-
"""
Placement Survey PPT 생성 - Summary RMS 차트
템플릿 기반, RMS Excel에서 데이터 읽어 차트 갱신
Usage: python gen_ppt.py
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

import re
import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.util import Pt
from lxml import etree

# ============================================================
# 고정 색상 매핑
# ============================================================
# JK 이유별 (범례 표시 순서 = dict 키 순서)
JK_REASON_COLORS = {
    '채용 공고':        RGBColor(0x4F, 0x81, 0xBD),  # blue
    '유용한 정보 제공':  RGBColor(0xC0, 0x50, 0x4D),  # red
    '이용 편의성':      RGBColor(0x9B, 0xBB, 0x59),  # green
    '브랜드/평판':      RGBColor(0x80, 0x64, 0xA2),  # purple
    '부가 서비스':      RGBColor(0x4B, 0xAC, 0xC6),  # teal
    '공공/제도/신뢰':   RGBColor(0x70, 0xAD, 0x47),  # teal
    '기타':            RGBColor(0xA5, 0xA5, 0xA5),  # grey
}
# AM 이유별 (범례 표시 순서 = dict 키 순서)
AM_REASON_COLORS = {
    '아르바이트 공고':        RGBColor(0x4F, 0x81, 0xBD),  # blue
    '유용한 정보 제공':       RGBColor(0xC0, 0x50, 0x4D),  # red
    '이용 편의성':           RGBColor(0x9B, 0xBB, 0x59),  # green
    '사용자 경험 및 인지도':   RGBColor(0x80, 0x64, 0xA2),  # purple
    '부가서비스 및 혜택':     RGBColor(0x4B, 0xAC, 0xC6),  # teal
    '위치 및 채용 특성':      RGBColor(0x70, 0xAD, 0x47),  # teal
    '그외 (기타)':           RGBColor(0xA5, 0xA5, 0xA5),  # grey
}
# Line chart 채널별
JK_LINE_COLORS = {
    'JK':  RGBColor(0x4F, 0x81, 0xBD),   # blue
    'SRI': RGBColor(0xC0, 0x50, 0x4D),   # red
    'RMB': RGBColor(0x9B, 0xBB, 0x59),   # green
    'WTD': RGBColor(0x80, 0x64, 0xA2),   # purple
}
AM_LINE_COLORS = {
    'AM': RGBColor(0x4F, 0x81, 0xBD),    # blue
    'AH': RGBColor(0xC0, 0x50, 0x4D),    # red
    'DG': RGBColor(0x9B, 0xBB, 0x59),    # green
}


def apply_bar_colors(chart, color_map):
    """Bar chart 시리즈에 이유별 고정 색상 + 데이터 레이블 % 포맷 적용."""
    _ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
    _ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    _ns_c14 = 'http://schemas.microsoft.com/office/drawing/2007/8/2/chart'
    _ns_mc = 'http://schemas.openxmlformats.org/markup-compatibility/2006'
    ns = {'c': _ns_c, 'a': _ns_a}

    # 차트 스타일 오버라이드 제거 (c14:style, c:style → 개별 색상 우선)
    cs = chart._chartSpace
    for style_el in cs.findall(f'{{{_ns_c}}}style'):
        cs.remove(style_el)
    for ac in cs.findall(f'{{{_ns_mc}}}AlternateContent'):
        # c14:style이 포함된 AlternateContent 제거
        if ac.find(f'.//{{{_ns_c14}}}style') is not None:
            cs.remove(ac)

    for s in chart.series:
        ser = s._element
        tx = ser.find('.//c:tx//c:v', ns)
        name = tx.text if tx is not None else ''

        # spPr 가져오거나 생성
        spPr = ser.find(f'{{{_ns_c}}}spPr')
        if spPr is None:
            spPr = etree.SubElement(ser, f'{{{_ns_c}}}spPr')

        # 기존 fill 제거
        for old in spPr.findall(f'{{{_ns_a}}}solidFill') + spPr.findall(f'{{{_ns_a}}}noFill'):
            spPr.remove(old)

        # solidFill을 spPr의 첫 번째 자식으로 삽입 (PowerPoint 표준 순서)
        if name and name.strip() and name != 'None' and name in color_map:
            rgb = color_map[name]
            sf = etree.Element(f'{{{_ns_a}}}solidFill')
            etree.SubElement(sf, f'{{{_ns_a}}}srgbClr', val=str(rgb))
            spPr.insert(0, sf)
        else:
            nf = etree.Element(f'{{{_ns_a}}}noFill')
            spPr.insert(0, nf)

        # 데이터 레이블 % 포맷 설정
        dLbls = ser.find(f'{{{_ns_c}}}dLbls')
        if dLbls is None:
            dLbls = etree.SubElement(ser, f'{{{_ns_c}}}dLbls')
        numFmt = dLbls.find(f'{{{_ns_c}}}numFmt')
        if numFmt is None:
            numFmt = etree.SubElement(dLbls, f'{{{_ns_c}}}numFmt')
        numFmt.set('formatCode', '0%')
        numFmt.set('sourceLinked', '0')


def apply_line_colors(chart, color_map):
    """Line chart 시리즈에 채널별 고정 색상 적용 (XML 직접 조작)."""
    _ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
    _ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns = {'c': _ns_c, 'a': _ns_a}

    for s in chart.series:
        ser = s._element
        tx = ser.find('.//c:tx//c:v', ns)
        name = tx.text if tx is not None else ''
        if name in color_map:
            rgb = color_map[name]
            spPr = ser.find(f'{{{_ns_c}}}spPr')
            if spPr is None:
                spPr = etree.SubElement(ser, f'{{{_ns_c}}}spPr')
            # ln 가져오거나 생성
            ln = spPr.find(f'{{{_ns_a}}}ln')
            if ln is None:
                ln = etree.SubElement(spPr, f'{{{_ns_a}}}ln', w='31750')  # ~2.5pt
            # 기존 fill 제거
            for old in ln.findall(f'{{{_ns_a}}}solidFill') + ln.findall(f'{{{_ns_a}}}noFill'):
                ln.remove(old)
            sf = etree.SubElement(ln, f'{{{_ns_a}}}solidFill')
            etree.SubElement(sf, f'{{{_ns_a}}}srgbClr', val=str(rgb))


# ============================================================
# Config
# ============================================================
BASE_DIR = Path(r"C:\Users\olivia408\projects\-dashboard\Placement survey")
TEMPLATE = BASE_DIR / "Placement Survey_25Q3.pptx"
JK_RMS = BASE_DIR / "JK 전체" / "25Q4_JK_RMS_v4.xlsx"
AM_RMS = BASE_DIR / "AM 전체" / "25Q4_AM_RMS_v3.xlsx"
OUTPUT = BASE_DIR / "Placement Survey_25Q4_output27.pptx"

# 3년치 분기 필터 (25Q4 기준 → 23Q1~25Q4)
JK_QUARTERS = ['23Q1','23Q2','23Q3','23Q4','24Q1','24Q2','24Q3','24Q4','25Q1','25Q2','25Q3','25Q4']
AM_QUARTERS = ['23.1Q','23.2Q','23.3Q','23.4Q','24.1Q','24.2Q','24.3Q','24.4Q','25.1Q','25.2Q','25.3Q','25.4Q']


# ============================================================
# Data Loading
# ============================================================
def load_rms_section(filepath, sheet, scope_label, channels, quarters_filter, section_index=0):
    """RMS 시트에서 특정 scope의 채널별 값을 로드
    section_index: 같은 라벨이 여러번 나올 때 (0=Adjusted, 1=Unadjusted)
    """
    df = pd.read_excel(filepath, sheet_name=sheet, header=None, engine='openpyxl')

    # quarter 헤더 행 찾기
    all_quarters = []
    for i, row in df.iterrows():
        qs = []
        for c in range(2, df.shape[1]):
            v = row.iloc[c]
            if v and 'Q' in str(v):
                qs.append(str(v).strip())
        if len(qs) >= 3:
            all_quarters = qs
            break

    if not all_quarters:
        return {}, []

    indices = [j for j, q in enumerate(all_quarters) if q in quarters_filter]
    filtered_q = [all_quarters[j] for j in indices]

    # scope_label 행 찾기 (N번째 occurrence)
    scope_row = None
    found_count = 0
    for i, row in df.iterrows():
        for val in row.values:
            if val and str(val).strip() == scope_label:
                if found_count == section_index:
                    scope_row = i
                    break
                found_count += 1
        if scope_row is not None:
            break

    if scope_row is None:
        print(f"  WARNING: '{scope_label}' (idx={section_index}) not found in {sheet}")
        return {}, filtered_q

    # 채널별 데이터 추출
    result = {}
    for ch in channels:
        for r in range(scope_row + 1, min(scope_row + len(channels) + 3, len(df))):
            cell = df.iloc[r, 1]
            if cell and str(cell).strip() == ch:
                vals = []
                for j in indices:
                    v = df.iloc[r, 2 + j]
                    try:
                        vals.append(round(float(v) * 100, 1))
                    except (ValueError, TypeError):
                        vals.append(0)
                result[ch] = vals
                break

    return result, filtered_q


def load_iar(filepath, sheet_list, target_channel, scope_label, quarters_filter, section_index=0):
    """인지/지원/재지원 각 시트에서 target_channel의 Share 값 로드
    section_index: 0=Adjusted블록, 1=Unadjusted블록
    """
    result = {}
    quarters = []

    for sheet, label in sheet_list:
        df = pd.read_excel(filepath, sheet_name=sheet, header=None, engine='openpyxl')

        if not quarters:
            for i, row in df.iterrows():
                qs = []
                for c in range(2, df.shape[1]):
                    v = row.iloc[c]
                    if v and 'Q' in str(v):
                        qs.append(str(v).strip())
                if len(qs) >= 3:
                    all_q = qs
                    break
            indices = [j for j, q in enumerate(all_q) if q in quarters_filter]
            quarters = [all_q[j] for j in indices]

        # scope 찾기 (N번째 occurrence)
        found = False
        found_count = 0
        for i, row in df.iterrows():
            if found:
                break
            for val in row.values:
                if val and scope_label in str(val):
                    if found_count == section_index:
                        for r in range(i + 1, min(i + 8, len(df))):
                            ch = df.iloc[r, 1]
                            if ch and str(ch).strip() == target_channel:
                                vals = []
                                for j in indices:
                                    v = df.iloc[r, 2 + j]
                                    try:
                                        vals.append(round(float(v) * 100, 1))
                                    except:
                                        vals.append(0)
                                result[label] = vals
                                found = True
                                break
                        break
                    found_count += 1

    return result, quarters


def iter_all_shapes(slide):
    """슬라이드의 모든 shape를 재귀적으로 순회 (Group 내부 포함)"""
    def _recurse(shapes):
        for shape in shapes:
            yield shape
            if shape.shape_type == 6:  # GROUP
                yield from _recurse(shape.shapes)
    yield from _recurse(slide.shapes)


# ============================================================
# Chart Update
# ============================================================
def update_chart(chart, categories, series_dict, zero_to_none=False):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_dict.items():
        if all(v == 0 for v in values):  # 전부 0이면 시리즈 제외
            continue
        if zero_to_none:
            values = [None if v == 0 else v for v in values]
        chart_data.add_series(name, values)
    chart.replace_data(chart_data)


def rename_jk(ch):
    return ch.replace('잡코리아','JK').replace('사람인','SRI').replace('원티드','WTD').replace('리멤버','RMB')


def rename_am(ch):
    return ch.replace('알바몬','AM').replace('알바천국','AH').replace('당근(알바)','DG').replace('온라인 Others','온라인')


# ============================================================
# Insight Generation
# ============================================================
INSIGHT_POS_Y = 5722070  # 시사점 텍스트 Y 위치


def find_insight_shape(slide):
    """시사점 텍스트 shape 찾기 (차트 아래, Y > 5000000, Group 내부 포함)"""
    for shape in iter_all_shapes(slide):
        if shape.has_text_frame and shape.top > 5000000:
            text = shape.text_frame.text.strip()
            if len(text) > 10:
                return shape
    return None


def set_insight_text(slide, text):
    """슬라이드 시사점 텍스트 교체"""
    shape = find_insight_shape(slide)
    if shape:
        for para in shape.text_frame.paragraphs:
            if para.text.strip():
                # 기존 폰트 스타일 유지하면서 텍스트만 교체
                for run in para.runs:
                    run.text = ''
                if para.runs:
                    para.runs[0].text = text
                else:
                    para.text = text
                break
        return True
    return False


def gen_insight_total(data, channels, brand, latest_q, prev_q):
    """Slide 4/11: 전체 채널 시사점 생성"""
    # 최신 분기 값
    vals = {ch: data[ch][-1] for ch in channels if ch in data}
    vals_prev = {ch: data[ch][-2] for ch in channels if ch in data and len(data[ch]) >= 2}

    # 1위, 2위 찾기
    online_chs = [ch for ch in channels if '오프라인' not in ch and '미분류' not in ch and 'Out' not in ch and 'Unpaid' not in ch and 'Paid' not in ch]
    sorted_ch = sorted(online_chs, key=lambda c: vals.get(c, 0), reverse=True)

    brand_val = vals.get(brand, 0)
    brand_prev = vals_prev.get(brand, 0)
    brand_dir = '상승' if brand_val > brand_prev else '하락' if brand_val < brand_prev else '유지'

    # 오프라인 비중
    offline_keys = [ch for ch in channels if '오프라인' in ch]
    offline_val = sum(vals.get(ch, 0) for ch in offline_keys)
    offline_prev = sum(vals_prev.get(ch, 0) for ch in offline_keys)
    offline_dir = '확대' if offline_val > offline_prev else '축소'

    # 기타 온라인
    etc_keys = [ch for ch in channels if '기타' in ch or 'Others' in ch]
    etc_val = sum(vals.get(ch, 0) for ch in etc_keys)

    parts = []
    if len(sorted_ch) >= 2:
        rival = sorted_ch[0] if sorted_ch[0] != brand else sorted_ch[1]
        rival_val = vals.get(rival, 0)
        parts.append(f"온라인 플랫폼 내 {rival}({rival_val:.1f}%) 외 주요 경쟁자는 보이지 않으며")

    if etc_val > 20:
        parts.append(f"기타 온라인({etc_val:.1f}%)과 오프라인({offline_val:.1f}%) 비중이 높은 상태")
    else:
        parts.append(f"오프라인 비중이 {offline_dir} 추세({offline_prev:.1f}%→{offline_val:.1f}%)")

    return ', '.join(parts)


def gen_insight_online(data, channels, brand, rival_name):
    """Slide 5/12: 온라인 경쟁 시사점 생성"""
    vals = {ch: data[ch][-1] for ch in channels if ch in data}
    vals_prev = {ch: data[ch][-2] for ch in channels if ch in data and len(data[ch]) >= 2}

    brand_val = vals.get(brand, 0)
    rival_val = vals.get(rival_name, 0)
    gap = brand_val - rival_val

    # 3위 이하
    others = [(ch, vals.get(ch, 0)) for ch in channels if ch != brand and ch != rival_name and ch in vals]
    others_sorted = sorted(others, key=lambda x: x[1], reverse=True)

    if abs(gap) < 3:
        compete = f"{rival_name}와 치열한 경쟁을 진행 중이며({brand} {brand_val:.1f}% vs {rival_name} {rival_val:.1f}%)"
    elif gap > 0:
        compete = f"{rival_name} 대비 {gap:.1f}%p 우위를 유지 중이며({brand} {brand_val:.1f}% vs {rival_name} {rival_val:.1f}%)"
    else:
        compete = f"{rival_name}에 {abs(gap):.1f}%p 열위 상태이며({brand} {brand_val:.1f}% vs {rival_name} {rival_val:.1f}%)"

    if others_sorted:
        minor = ', '.join(f"{ch}({v:.1f}%)" for ch, v in others_sorted)
        threat = f"{minor}는 위협할만한 경쟁자로 인식되지 않음"
    else:
        threat = ""

    return f"{compete}, {threat}" if threat else compete


def gen_insight_iar(iar_data, brand_name):
    """Slide 6/13: 인지/지원/재지원 갭 시사점 생성"""
    apply_vals = iar_data.get('지원', [])
    aware_vals = iar_data.get('인지', [])
    reuse_vals = iar_data.get('재지원', [])

    if not apply_vals or not reuse_vals:
        return ""

    # 최근 값
    apply_now = apply_vals[-1]
    aware_now = aware_vals[-1] if aware_vals else 0
    reuse_now = reuse_vals[-1]

    # 지원 vs 인지 갭
    aware_gap = apply_now - aware_now
    # 지원 vs 재지원 갭
    reuse_gap = apply_now - reuse_now

    parts = []

    # 인지 vs 지원
    if abs(aware_gap) < 2:
        parts.append(f"인지({aware_now:.1f}%)와 지원({apply_now:.1f}%) RMS가 유사한 수준")
    elif aware_gap > 0:
        parts.append(f"지원({apply_now:.1f}%)이 인지({aware_now:.1f}%) 대비 {aware_gap:.1f}%p 높음")
    else:
        parts.append(f"인지({aware_now:.1f}%)가 지원({apply_now:.1f}%) 대비 {abs(aware_gap):.1f}%p 높음")

    # 재지원 갭 및 추세
    if len(apply_vals) >= 4 and len(reuse_vals) >= 4:
        gaps = [apply_vals[i] - reuse_vals[i] for i in range(-4, 0)]
        gap_trend = gaps[-1] - gaps[0]
        if reuse_gap > 1:
            if gap_trend < -1:
                parts.append(f"재지원({reuse_now:.1f}%)이 지원 대비 {reuse_gap:.1f}%p 낮으나, 격차가 감소하는 추이")
            elif gap_trend > 1:
                parts.append(f"재지원({reuse_now:.1f}%)이 지원 대비 {reuse_gap:.1f}%p 낮으며, 격차가 확대되는 추이")
            else:
                parts.append(f"재지원({reuse_now:.1f}%)이 지원 대비 {reuse_gap:.1f}%p 낮게 유지 중")
        elif reuse_gap < -1:
            parts.append(f"재지원({reuse_now:.1f}%)이 지원 대비 {abs(reuse_gap):.1f}%p 높아 충성도 확보 중")
        else:
            parts.append(f"재지원({reuse_now:.1f}%)과 지원이 유사한 수준")
    else:
        if reuse_gap > 1:
            parts.append(f"재지원({reuse_now:.1f}%)이 지원 대비 {reuse_gap:.1f}%p 낮게 형성")
        else:
            parts.append(f"재지원({reuse_now:.1f}%)과 지원이 유사한 수준")

    return ', '.join(parts)


# ============================================================
# Seg Scatter Chart (Slide 7-10, 14-17)
# ============================================================
def load_seg_graph_data(filepath, sheet_name, target_channel, seg_dims):
    """Seg 그래프 시트에서 각 dimension별 (seg_name, 1yr, 3yr) 로드
    가로 배치: 각 dim이 col 0, 6, 12, 18 등에 위치
    Returns: {dim_name: [(seg_name, x_1yr, y_3yr), ...]}
    """
    df = pd.read_excel(filepath, sheet_name=sheet_name, header=None, engine='openpyxl')
    result = {}

    # Row 3에서 dimension header 위치 찾기
    dim_cols = {}
    for c in range(df.shape[1]):
        v = df.iloc[3, c] if 3 < len(df) else None
        if pd.notna(v) and str(v).startswith('['):
            dim_name = str(v).strip().strip('[]')
            if dim_name in seg_dims:
                dim_cols[dim_name] = c

    for dim_name, col_start in dim_cols.items():
        result[dim_name] = []
        # col_start = seg name, col_start+1 = channel, col_start+2 = 1yr, col_start+3 = 3yr
        current_seg = None
        for i in range(5, len(df)):
            seg_val = df.iloc[i, col_start]
            ch_val = df.iloc[i, col_start + 1]

            # Seg name row (has "(N" in channel col)
            if pd.notna(seg_val) and pd.notna(ch_val) and '(N' in str(ch_val):
                current_seg = str(seg_val).strip()
                continue

            # Target channel row
            if pd.notna(ch_val) and str(ch_val).strip() == target_channel and current_seg:
                v_1yr = df.iloc[i, col_start + 2]
                v_3yr = df.iloc[i, col_start + 3]
                try:
                    x = float(v_1yr)
                    y = float(v_3yr)
                    if x > 0 or y > 0:  # 유효한 데이터만
                        result[dim_name].append((current_seg, x, y))
                except (ValueError, TypeError):
                    pass

    return result


def update_scatter_chart(chart, points):
    """XY Scatter 차트 데이터 교체: points = [(x, y), ...]"""
    chart_data = XyChartData()
    series = chart_data.add_series('Present')
    for x, y in points:
        series.add_data_point(x, y)
    chart.replace_data(chart_data)


# ============================================================
# Double Click Logic
# ============================================================
def load_seg_cut_timeseries(filepath, sheet_name, seg_dims_cols):
    """Seg Cut 시트에서 seg별 4사 RMS 시계열 로드
    seg_dims_cols: {dim_name: col_start}
    Returns: {dim_name: {seg_name: {channel: [val_per_quarter]}}}, [quarters]
    """
    df = pd.read_excel(filepath, sheet_name=sheet_name, header=None, engine='openpyxl')

    # Find dimension column starts from row 2
    dim_cols = {}
    for c in range(df.shape[1]):
        v = df.iloc[2, c]
        if pd.notna(v) and str(v).startswith('['):
            name = str(v).strip().strip('[]')
            dim_cols[name] = c

    # Quarters from row 3 of first dimension
    quarters = []
    first_col = list(dim_cols.values())[0] if dim_cols else 0
    for c in range(first_col + 2, first_col + 20):
        v = df.iloc[3, c]
        if pd.notna(v) and 'Q' in str(v):
            quarters.append(str(v).strip())

    # 4사 채널 목록 (JK vs AM)
    jk_4sa = {'잡코리아', '사람인', '원티드', '리멤버'}
    am_3sa = {'알바몬', '알바천국', '당근(알바)'}

    result = {}
    for dim_name, col_s in dim_cols.items():
        result[dim_name] = {}
        current_seg = None
        read_count = 0  # 현재 seg에서 읽은 채널 수

        for i in range(4, len(df)):
            v0 = df.iloc[i, col_s]
            v1 = df.iloc[i, col_s + 1]

            # New seg name (has "(N)" in channel col)
            if pd.notna(v0) and pd.notna(v1) and '(N)' in str(v1):
                current_seg = str(v0).strip()
                read_count = 0
                result[dim_name][current_seg] = {}
                continue

            # Channel data — only first block (4사 Share)
            if pd.notna(v1) and current_seg:
                ch = str(v1).strip()
                # 이미 읽은 채널이면 두 번째 블록 → skip
                if ch in result[dim_name][current_seg]:
                    continue
                if ch in jk_4sa or ch in am_3sa:
                    vals = []
                    for q_idx in range(len(quarters)):
                        v = df.iloc[i, col_s + 2 + q_idx]
                        try:
                            vals.append(round(float(v) * 100, 1))
                        except:
                            vals.append(0)
                    result[dim_name][current_seg][ch] = vals

    return result, quarters


def load_seg_reason_all(filepath, sheet_name):
    """Seg×채널 이유 시트를 한 번에 로드.
    Returns: {dim_name: {seg_val: {channel: {reason: share}}}}
    """
    df = pd.read_excel(filepath, sheet_name=sheet_name, header=None, engine='openpyxl')
    result = {}
    current_dim = None
    reason_cats = []

    for i in range(len(df)):
        v0 = df.iloc[i, 0]

        # Dimension header: [산업], [직무] etc
        if pd.notna(v0) and str(v0).startswith('[') and str(v0).endswith(']'):
            current_dim = str(v0).strip().strip('[]')
            result[current_dim] = {}
            continue

        # Header row: Seg, Channel, N, reason1, reason2, ...
        if pd.notna(v0) and str(v0).strip() == 'Seg':
            reason_cats = []
            for c in range(3, df.shape[1]):
                v = df.iloc[i, c]
                if pd.notna(v):
                    reason_cats.append(str(v).strip())
            continue

        # Data row
        if current_dim and pd.notna(df.iloc[i, 1]):
            seg_val = str(df.iloc[i, 0]).strip() if pd.notna(df.iloc[i, 0]) else None
            ch = str(df.iloc[i, 1]).strip()
            if ch == 'Channel' or not ch:
                continue

            # Find seg name (could be empty if continuation)
            if seg_val and seg_val != 'nan':
                last_seg = seg_val
            else:
                seg_val = last_seg if 'last_seg' in dir() else None

            if not seg_val or seg_val == 'nan':
                continue

            if seg_val not in result[current_dim]:
                result[current_dim][seg_val] = {}

            shares = {}
            for j, cat in enumerate(reason_cats):
                v = df.iloc[i, 3 + j]
                try:
                    shares[cat] = float(v)
                except:
                    shares[cat] = 0
            result[current_dim][seg_val][ch] = shares

    return result


def get_reason_top3(reason_shares):
    """이유 비중 dict에서 Top3 추출. Returns: [(reason, share), ...]"""
    if not reason_shares:
        return []
    sorted_reasons = sorted(reason_shares.items(), key=lambda x: x[1], reverse=True)
    return [(r, s) for r, s in sorted_reasons[:3] if s > 0]


def update_double_click_line(chart, seg_data, channels, quarters, q_filter):
    """Double Click line chart 업데이트"""
    indices = [i for i, q in enumerate(quarters) if q in q_filter]
    filtered_q = [quarters[i] for i in indices]

    chart_data = CategoryChartData()
    chart_data.categories = filtered_q

    for ch in channels:
        vals = seg_data.get(ch, [])
        filtered_vals = [vals[i] if i < len(vals) else 0 for i in indices]
        chart_data.add_series(ch, filtered_vals)

    chart.replace_data(chart_data)


def update_double_click_bar(chart, reason_top3):
    """Double Click bar chart 업데이트: 이유 Top3"""
    if not reason_top3:
        return
    chart_data = CategoryChartData()
    # Bar chart has 1 category (channel name) and 3 series (reasons)
    cats = list(chart.plots[0].categories) if chart.plots else ['']
    chart_data.categories = cats

    for reason, pct in reason_top3:
        chart_data.add_series(reason, [round(pct, 1)])

    chart.replace_data(chart_data)


def gen_double_click_insight(seg_name, seg_data, brand, rival, quarters_used):
    """Double Click 시사점 생성"""
    brand_vals = seg_data.get(brand, [])
    rival_vals = seg_data.get(rival, [])

    if not brand_vals or not rival_vals:
        return ""

    b_now = brand_vals[-1]
    r_now = rival_vals[-1]
    gap = b_now - r_now

    # Trend (last 4 quarters)
    recent = brand_vals[-4:] if len(brand_vals) >= 4 else brand_vals
    if len(recent) >= 2:
        trend = recent[-1] - recent[0]
        trend_txt = '상승' if trend > 2 else '하락' if trend < -2 else '유지'
    else:
        trend_txt = '확인 필요'

    short_brand = brand.replace('잡코리아', 'JK').replace('사람인', 'SRI').replace('알바몬', 'AM').replace('알바천국', 'AH')
    short_rival = rival.replace('잡코리아', 'JK').replace('사람인', 'SRI').replace('알바몬', 'AM').replace('알바천국', 'AH')

    if gap > 10:
        compete = f"{short_rival} 대비 {gap:.0f}%p 우위를 유지하며 RMS {trend_txt} 추세"
    elif gap > 0:
        compete = f"{short_rival}와 소폭 우위({short_brand} {b_now:.1f}% vs {short_rival} {r_now:.1f}%), RMS {trend_txt} 추세"
    elif gap > -10:
        compete = f"{short_rival}에 소폭 열위({short_brand} {b_now:.1f}% vs {short_rival} {r_now:.1f}%), RMS {trend_txt} 추세"
    else:
        compete = f"{short_rival}에 {abs(gap):.0f}%p 열위 상태이며 RMS {trend_txt} 추세"

    return compete


def gen_seg_insight(points, seg_names, brand_name):
    """Seg scatter 기반 시사점: 4분면 분류 후 요약"""
    if not points:
        return ""

    avg_x = sum(p[0] for p in points) / len(points)
    avg_y = sum(p[1] for p in points) / len(points)

    q1 = []  # 과거↑ 최근↑ (강점)
    q2 = []  # 과거↑ 최근↓ (악화)
    q3 = []  # 과거↓ 최근↓ (약점)
    q4 = []  # 과거↓ 최근↑ (개선)

    for name, x, y in zip(seg_names, [p[0] for p in points], [p[1] for p in points]):
        if x >= avg_x and y >= avg_y:
            q1.append(name)
        elif x < avg_x and y >= avg_y:
            q2.append(name)
        elif x < avg_x and y < avg_y:
            q3.append(name)
        else:
            q4.append(name)

    parts = []
    if q2:
        parts.append(f"최우선 개선 필요: {', '.join(q2)}")
    if q3:
        parts.append(f"지속 약세: {', '.join(q3)}")
    if q4:
        parts.append(f"개선 추세: {', '.join(q4)}")
    if q1:
        parts.append(f"강점 유지: {', '.join(q1)}")

    return ' / '.join(parts)


# ============================================================
# Main
# ============================================================
def main():
    print("=" * 50)
    print("Placement Survey PPT 생성")
    print("=" * 50)

    prs = Presentation(str(TEMPLATE))
    print(f"[1] Template: {TEMPLATE.name}")

    # ========================
    # JK (상용직)
    # ========================
    print(f"\n[2] JK RMS loading...")

    # Slide 4: 전체 채널 Unadjusted
    jk_all_ch = ['잡코리아','사람인','원티드','리멤버','기타 온라인','오프라인','미분류']
    jk_total, q_jk = load_rms_section(
        JK_RMS, '지원 RMS', '전체 (Online+Offline)', jk_all_ch, JK_QUARTERS,
        section_index=1)  # Unadjusted
    print(f"  전체 Unadjusted: {len(jk_total)} ch, {len(q_jk)} Q")

    # Slide 5: 4사 Share (Unadjusted 기준, 4사 합=100%)
    jk_4 = ['잡코리아','사람인','원티드','리멤버']
    jk_online, _ = load_rms_section(
        JK_RMS, '지원 RMS', 'JK SRI WTD RMB Share', jk_4, JK_QUARTERS,
        section_index=1)  # Unadjusted
    print(f"  4사 Unadj Share: {len(jk_online)} ch")

    # Slide 6: 인지/지원/재지원 (Unadjusted 4사 Share JK)
    jk_iar, q_iar = load_iar(
        JK_RMS,
        [('지원 RMS','지원'), ('인지 RMS','인지'), ('재지원 RMS','재지원')],
        '잡코리아', 'JK SRI WTD RMB Share', JK_QUARTERS,
        section_index=1)  # Unadjusted
    print(f"  인지/지원/재지원: {list(jk_iar.keys())}")

    # ========================
    # AM (일용직)
    # ========================
    print(f"\n[3] AM RMS loading...")

    am_all_ch = ['알바몬','알바천국','당근(알바)','온라인 Others','오프라인 Paid','오프라인 Unpaid']
    am_total, q_am = load_rms_section(
        AM_RMS, '지원 RMS', '전체 (Online+Offline)', am_all_ch, AM_QUARTERS,
        section_index=1)  # Unadjusted
    print(f"  전체 Unadjusted: {len(am_total)} ch, {len(q_am)} Q")

    # Slide 12: 3사 Share (Unadjusted 기준, 3사 합=100%)
    am_3 = ['알바몬','알바천국','당근(알바)']
    am_online, _ = load_rms_section(
        AM_RMS, '지원 RMS', 'AM AH DG Share', am_3, AM_QUARTERS,
        section_index=1)  # Unadjusted
    print(f"  3사 Unadj Share: {len(am_online)} ch")

    am_iar, q_am_iar = load_iar(
        AM_RMS,
        [('지원 RMS','지원'), ('인지 RMS','인지'), ('재지원 RMS','재지원')],
        '알바몬', 'AM AH DG Share', AM_QUARTERS,
        section_index=1)  # Unadjusted
    print(f"  인지/지원/재지원: {list(am_iar.keys())}")

    # ========================
    # 차트 업데이트 (25Q3 템플릿 기준 슬라이드 인덱스)
    # ========================
    print(f"\n[4] Updating charts...")

    # Slide 4 (idx 3): JK 전체 채널 Share
    if jk_total:
        for shape in prs.slides[3].shapes:
            if shape.has_chart:
                s = {rename_jk(ch): jk_total[ch] for ch in jk_all_ch if ch in jk_total}
                update_chart(shape.chart, q_jk, s)
                print(f"  Slide 4: JK 전체 채널 Share ✓")

    # Slide 5 (idx 4): JK 4사 Share
    if jk_online:
        for shape in prs.slides[4].shapes:
            if shape.has_chart:
                s = {rename_jk(ch): jk_online[ch] for ch in jk_4 if ch in jk_online}
                update_chart(shape.chart, q_jk, s)
                print(f"  Slide 5: JK 4사 Share ✓")

    # Slide 6 (idx 5): JK 인지/지원/재지원
    if jk_iar:
        for shape in prs.slides[5].shapes:
            if shape.has_chart:
                update_chart(shape.chart, q_iar, jk_iar)
                print(f"  Slide 6: JK 인지/지원/재지원 ✓")

    # Slide 11 (idx 10): AM 전체 채널 Share
    if am_total:
        for shape in prs.slides[10].shapes:
            if shape.has_chart:
                s = {rename_am(ch): am_total[ch] for ch in am_all_ch if ch in am_total}
                update_chart(shape.chart, q_am, s)
                print(f"  Slide 11: AM 전체 채널 Share ✓")

    # Slide 12 (idx 11): AM 3사 Share
    if am_online:
        for shape in prs.slides[11].shapes:
            if shape.has_chart:
                s = {rename_am(ch): am_online[ch] for ch in am_3 if ch in am_online}
                update_chart(shape.chart, q_am, s)
                print(f"  Slide 12: AM 3사 Share ✓")

    # Slide 13 (idx 12): AM 인지/지원/재지원
    if am_iar:
        for shape in prs.slides[12].shapes:
            if shape.has_chart:
                update_chart(shape.chart, q_am_iar, am_iar)
                print(f"  Slide 13: AM 인지/지원/재지원 ✓")

    # Slides 7-10, 14-17 (Seg Scatter): 수작업 — 건드리지 않음
    _SKIP_SEG_SCATTER = True

    # --- JK/AM Seg Scatter (Slides 7-10, 14-17) — 비활성화 ---
    if not _SKIP_SEG_SCATTER:
        pass  # scatter chart + quadrant 코드 (필요시 _SKIP_SEG_SCATTER = False로 활성화)

    # ========================
    # 표지 업데이트 (Slide 1)
    # ========================
    latest_jk_q = JK_QUARTERS[-1]  # e.g. "25Q4"
    q_display = latest_jk_q.replace('Q', ' ').replace('25', '25 ')  # "25 4"
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    # "25 3Q" → "25 4Q" 등 분기 텍스트 교체
                    if 'Q' in run.text and any(c.isdigit() for c in run.text):
                        old = run.text
                        # 25 3Q → 25 4Q, 25Q3 → 25Q4 등
                        q_num = latest_jk_q[-1]  # "25Q4" → "4"
                        run.text = re.sub(r'25\s*\d\s*Q', f'25 {q_num}Q', run.text)
                        run.text = re.sub(r'25Q\d', latest_jk_q, run.text)
                        if old != run.text:
                            print(f"  Slide 1: \"{old}\" → \"{run.text}\"")
    print(f"  표지 → {latest_jk_q}")

    # ========================
    # 재지원 Double Click (Slide 19, 57) — Bar chart 업데이트
    # ========================
    # JK 재지원 이유 로드
    jk_reuse_reasons = load_seg_reason_all(JK_RMS, '재지원 Seg×채널 이유') if '재지원 Seg×채널 이유' in pd.ExcelFile(JK_RMS, engine='openpyxl').sheet_names else {}
    am_reuse_reasons = load_seg_reason_all(AM_RMS, '재지원 Seg×채널 이유') if '재지원 Seg×채널 이유' in pd.ExcelFile(AM_RMS, engine='openpyxl').sheet_names else {}

    # Slide 19: JK 재지원 — 전체 채널별 이유 (채널 선택 이유 시트에서 로드)
    jk_reuse_reason_all = {}
    try:
        df_rr = pd.read_excel(JK_RMS, sheet_name='재지원 채널 선택 이유', header=None, engine='openpyxl')
        current_ch = None
        for i in range(len(df_rr)):
            v0 = df_rr.iloc[i, 0]
            v1 = df_rr.iloc[i, 1]
            if pd.notna(v0) and str(v0).strip() not in ('', '재지원 채널 선택 이유', '가중치: 1순위×3, 2순위×2, 3순위×1'):
                current_ch = str(v0).strip()
                jk_reuse_reason_all[current_ch] = {}
            if current_ch and pd.notna(v1) and str(v1).strip() != '이유 분류':
                cat = str(v1).strip()
                # 마지막 비중 컬럼 (최신Q)
                for c in range(df_rr.shape[1] - 1, 2, -1):
                    v = df_rr.iloc[i, c]
                    if pd.notna(v):
                        try:
                            jk_reuse_reason_all[current_ch][cat] = float(v)
                        except:
                            pass
                        break
    except:
        pass

    jk_reuse_display = {'잡코리아': 'JK', '사람인': 'SRI', '원티드': 'WTD', '리멤버': 'RMB'}
    jk_reuse_rev = {v: k for k, v in jk_reuse_display.items()}
    ns_chart = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
    if jk_reuse_reason_all:
        slide = prs.slides[18]  # Slide 19
        for shape in iter_all_shapes(slide):
            if shape.has_chart and shape.chart.chart_type == 51:
                cats = list(shape.chart.plots[0].categories)
                ch_display = cats[0] if cats else ''
                ch_raw = jk_reuse_rev.get(ch_display, ch_display)
                if ch_raw in jk_reuse_reason_all:
                    top3 = get_reason_top3(jk_reuse_reason_all[ch_raw])
                    if top3:
                        orig_count = len(list(shape.chart.series))
                        chart_data = CategoryChartData()
                        chart_data.categories = cats
                        for reason, share in top3:
                            chart_data.add_series(reason, [round(share, 2)])
                        for _ in range(max(orig_count, 3) - len(top3)):
                            chart_data.add_series(' ', [0])
                        shape.chart.replace_data(chart_data)
                        apply_bar_colors(shape.chart, JK_REASON_COLORS)
            # "(25Q3)" → "(25Q4)" 텍스트
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if '25Q' in run.text:
                            run.text = re.sub(r'25Q\d', latest_jk_q, run.text)
        print(f"  Slide 19: JK 재지원 이유 ✓")

    # Slide 57: AM 재지원
    am_reuse_reason_all = {}
    try:
        df_rr = pd.read_excel(AM_RMS, sheet_name='재지원 채널 선택 이유', header=None, engine='openpyxl')
        current_ch = None
        for i in range(len(df_rr)):
            v0 = df_rr.iloc[i, 0]
            v1 = df_rr.iloc[i, 1]
            if pd.notna(v0) and str(v0).strip() not in ('', '재지원 채널 선택 이유', '가중치: 1순위×3, 2순위×2, 3순위×1'):
                current_ch = str(v0).strip()
                am_reuse_reason_all[current_ch] = {}
            if current_ch and pd.notna(v1) and str(v1).strip() != '이유 분류':
                cat = str(v1).strip()
                for c in range(df_rr.shape[1] - 1, 2, -1):
                    v = df_rr.iloc[i, c]
                    if pd.notna(v):
                        try:
                            am_reuse_reason_all[current_ch][cat] = float(v)
                        except:
                            pass
                        break
    except:
        pass

    am_reuse_display = {'알바몬': 'AM', '알바천국': 'AH', '당근(알바)': 'DG'}
    am_reuse_rev = {v: k for k, v in am_reuse_display.items()}
    latest_am_q = AM_QUARTERS[-1]  # e.g. "25.4Q"
    if am_reuse_reason_all:
        slide = prs.slides[56]  # Slide 57
        for shape in iter_all_shapes(slide):
            if shape.has_chart and shape.chart.chart_type == 51:
                cats = list(shape.chart.plots[0].categories)
                ch_display = cats[0] if cats else ''
                ch_raw = am_reuse_rev.get(ch_display, ch_display)
                if ch_raw in am_reuse_reason_all:
                    top3 = get_reason_top3(am_reuse_reason_all[ch_raw])
                    if top3:
                        orig_count = len(list(shape.chart.series))
                        chart_data = CategoryChartData()
                        chart_data.categories = cats
                        for reason, share in top3:
                            chart_data.add_series(reason, [round(share, 2)])
                        for _ in range(max(orig_count, 3) - len(top3)):
                            chart_data.add_series(' ', [0])
                        shape.chart.replace_data(chart_data)
                        apply_bar_colors(shape.chart, AM_REASON_COLORS)
            # "(25Q3)" → "(25.4Q)" 텍스트
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if '25Q' in run.text or '25.' in run.text:
                            run.text = re.sub(r'25[Q.]\d[Q]?', latest_am_q, run.text)
        print(f"  Slide 57: AM 재지원 이유 ✓")

    # ========================
    # Double Click 슬라이드 업데이트
    # ========================
    print(f"\n[6] Double Click slides...")

    # JK Seg Cut 시계열 로드
    jk_seg_ts, jk_seg_quarters = load_seg_cut_timeseries(JK_RMS, '지원 Seg Cut', {})
    jk_channels_4sa = ['잡코리아', '사람인', '원티드', '리멤버']
    jk_ch_display = {'잡코리아': 'JK', '사람인': 'SRI', '원티드': 'WTD', '리멤버': 'RMB'}
    jk_display_to_raw = {v: k for k, v in jk_ch_display.items()}

    # AM Seg Cut 시계열 로드
    am_seg_ts, am_seg_quarters = load_seg_cut_timeseries(AM_RMS, '지원 Seg Cut', {})
    am_channels_3sa = ['알바몬', '알바천국', '당근(알바)']
    am_ch_display = {'알바몬': 'AM', '알바천국': 'AH', '당근(알바)': 'DG'}
    am_display_to_raw = {v: k for k, v in am_ch_display.items()}

    # Seg×채널 이유 비중 로드
    jk_seg_reasons = load_seg_reason_all(JK_RMS, '지원 Seg×채널 이유')
    am_seg_reasons = load_seg_reason_all(AM_RMS, '지원 Seg×채널 이유')

    def find_data_range(seg_data, all_quarters, max_quarters):
        """seg_data에서 실제 데이터 있는 기간 찾기 (최대 max_quarters)"""
        # 어떤 채널이든 데이터가 있는 첫 분기 ~ 마지막 분기
        first_idx = len(all_quarters)
        last_idx = 0
        for ch, vals in seg_data.items():
            for i, v in enumerate(vals):
                if v and v > 0:
                    first_idx = min(first_idx, i)
                    last_idx = max(last_idx, i)
        if first_idx > last_idx:
            return [], []
        # 3년(12Q) 제한: 뒤에서부터 최대 max_quarters
        start = max(first_idx, last_idx - max_quarters + 1)
        indices = list(range(start, last_idx + 1))
        return indices, [all_quarters[i] for i in indices]

    def update_dc_slide(slide, seg_data, seg_name, dim_short,
                        all_quarters, ch_list, ch_display_map, ch_reverse_map,
                        reason_seg, main_channels, threshold_channels,
                        brand, rival, line_colors=None, bar_colors=None, max_q=12):
        """Double Click 슬라이드 업데이트 (line + bar + subtitle + insight)"""
        indices, filtered_q = find_data_range(seg_data, all_quarters, max_q)
        if not filtered_q:
            return False

        period_str = f"({filtered_q[0]}~{filtered_q[-1]})"

        # --- Line chart --- (Group 내부도 탐색)
        # 원본 시리즈 순서 보존 (색상 유지 위해)
        for shape in iter_all_shapes(slide):
            if shape.has_chart and shape.chart.chart_type == 65:  # LINE_MARKERS
                # 원본 시리즈 순서 읽기
                orig_order = []
                for s in shape.chart.series:
                    ser = s._element
                    _ns = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
                    tx_el = ser.find('.//{%s}tx//{%s}v' % (_ns['c'], _ns['c']))
                    orig_order.append(tx_el.text if tx_el is not None else '?')

                chart_data = CategoryChartData()
                chart_data.categories = filtered_q

                for ch_disp in orig_order:
                    ch_raw = ch_reverse_map.get(ch_disp, ch_disp)
                    vals = seg_data.get(ch_raw, [])
                    fv = [vals[i] if i < len(vals) else 0 for i in indices]

                    show = True
                    if not any(v != 0 for v in fv):
                        show = False
                    if ch_raw in threshold_channels:
                        latest_val = fv[-1] if fv else 0
                        if (latest_val or 0) < 2.0:
                            show = False

                    if show:
                        fv_clean = [None if v == 0 else v for v in fv]
                        chart_data.add_series(ch_disp, fv_clean)

                shape.chart.replace_data(chart_data)
                if line_colors:
                    apply_line_colors(shape.chart, line_colors)
                break

        # --- Subtitle 업데이트 (기간 텍스트 + 이유 Q 텍스트) ---
        latest_q_str = filtered_q[-1]  # "25Q4" 또는 "25.4Q"
        for shape in iter_all_shapes(slide):
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text
                    # "RMS 추이 (23Q1~25Q3)" → "(23Q1~25Q4)"
                    if 'RMS 추이' in t:
                        for run in para.runs:
                            if '(' in run.text and 'Q' in run.text:
                                run.text = re.sub(r'\([^)]*Q[^)]*\)', period_str, run.text)
                    # "채널 선택 이유 Top 3 (25Q3)" → "(25Q4)"
                    if 'Top 3' in t or 'Top3' in t:
                        for run in para.runs:
                            if '(' in run.text and 'Q' in run.text:
                                run.text = re.sub(r'\([^)]*Q[^)]*\)', f'({latest_q_str})', run.text)

        # --- Bar charts (이유 Top3) ---
        # 활성 채널 결정: main + threshold 조건 충족
        show_ch_raw = list(main_channels)  # 순서 유지
        for ch in threshold_channels:
            vals = seg_data.get(ch, [])
            latest_val = vals[indices[-1]] if indices and indices[-1] < len(vals) else 0
            if (latest_val or 0) >= 2.0:
                show_ch_raw.append(ch)
        show_ch_disp = [ch_display_map.get(ch, ch) for ch in show_ch_raw]

        # 이 Seg의 모든 활성 채널 Top3에 등장하는 이유 수집 (고정 순서)
        fixed_order = list(bar_colors.keys()) if bar_colors else []
        all_top3_reasons = set()
        for ch_raw in show_ch_raw:
            for r, _ in get_reason_top3(reason_seg.get(ch_raw, {})):
                all_top3_reasons.add(r)
        # 범례에 표시할 이유 = 고정 순서 중 실제 Top3에 등장한 것만
        legend_reasons = [r for r in fixed_order if r in all_top3_reasons]

        if reason_seg:
            # Bar chart shapes 수집 (위치 left 순으로 정렬)
            bar_shapes = []
            for shape in iter_all_shapes(slide):
                if shape.has_chart and shape.chart.chart_type == 51:
                    bar_shapes.append(shape)
            bar_shapes.sort(key=lambda s: s.left)

            if len(bar_shapes) == 1 and len(show_ch_raw) > 1:
                # Grouped bar (한 차트에 여러 채널)
                shape = bar_shapes[0]
                try:
                    chart_data = CategoryChartData()
                    chart_data.categories = show_ch_disp
                    for reason in legend_reasons:
                        vals = []
                        for ch_raw in show_ch_raw:
                            ch_reasons = reason_seg.get(ch_raw, {})
                            ch_top3 = set(r for r, _ in get_reason_top3(ch_reasons))
                            if reason in ch_top3:
                                vals.append(round(ch_reasons.get(reason, 0), 2))
                            else:
                                vals.append(None)
                        chart_data.add_series(reason, vals)
                    shape.chart.replace_data(chart_data)
                    if bar_colors:
                        apply_bar_colors(shape.chart, bar_colors)
                except Exception:
                    pass
            else:
                # 개별 채널별 bar chart (left 순 매핑)
                for i, shape in enumerate(bar_shapes):
                    try:
                        if i < len(show_ch_raw):
                            ch_raw = show_ch_raw[i]
                            ch_disp = show_ch_disp[i]
                            ch_reasons = reason_seg.get(ch_raw, {})
                            top3 = get_reason_top3(ch_reasons)
                            top3_names = set(r for r, _ in top3)

                            chart_data = CategoryChartData()
                            chart_data.categories = [ch_disp]
                            # 고정 순서, Top3만 (0값 제외)
                            for reason in fixed_order:
                                if reason in top3_names:
                                    chart_data.add_series(reason, [round(ch_reasons.get(reason, 0), 2)])
                            shape.chart.replace_data(chart_data)
                            if bar_colors:
                                apply_bar_colors(shape.chart, bar_colors)
                        else:
                            # 초과 bar chart 비우기
                            chart_data = CategoryChartData()
                            chart_data.categories = ['']
                            chart_data.add_series(' ', [0])
                            shape.chart.replace_data(chart_data)
                    except Exception:
                        pass

        # --- 범례 업데이트 (사각형+텍스트 쌍) ---
        if bar_colors:
            # 범례 영역 shape 찾기 (x > 2800000, y > 4200000)
            legend_pairs = []
            rects = {}
            texts = {}
            for shape in iter_all_shapes(slide):
                x, y = shape.left, shape.top
                if x > 2800000 and y > 4200000 and y < 5800000:
                    if shape.shape_type == 1 and shape.width < 200000:
                        rects[y] = shape
                    elif shape.shape_type == 17 and shape.has_text_frame:
                        texts[y] = shape

            for ry, rect in sorted(rects.items()):
                closest_ty = min(texts.keys(), key=lambda ty: abs(ty - ry), default=None)
                if closest_ty is not None and abs(closest_ty - ry) < 200000:
                    legend_pairs.append((rect, texts.pop(closest_ty)))

            for i, (rect, textbox) in enumerate(legend_pairs):
                if i < len(legend_reasons):
                    reason_name = legend_reasons[i]
                    for para in textbox.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = reason_name
                    if reason_name in bar_colors:
                        rect.fill.solid()
                        rect.fill.fore_color.rgb = bar_colors[reason_name]
                else:
                    # 초과 범례 숨기기
                    for para in textbox.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = ''
                    # 사각형도 투명화
                    try:
                        rect.fill.background()
                    except Exception:
                        pass

        # --- Insight ---
        txt = gen_double_click_insight(seg_name, seg_data, brand, rival, filtered_q)
        set_insight_text(slide, txt)
        return True

    # JK Double Click 공통 설정
    jk_dim_map = {'소득': '소득수준', '산업': '산업', '직무': '직무', '지역': '지역'}
    jk_main_ch = ['잡코리아', '사람인']
    jk_threshold_ch = ['원티드', '리멤버']

    def run_jk_dc(slide_range, label):
        for idx in slide_range:
            slide = prs.slides[idx]
            title = ''
            for shape in slide.shapes:
                if shape.has_text_frame and 'Double Click' in shape.text_frame.text:
                    title = shape.text_frame.text.strip()
                    break
            if not title or '재지원' in title:
                continue
            m = re.match(r'\[(.+?)\]\s*(.+?)\s*Double Click', title)
            if not m:
                continue
            dim_short, seg_name = m.group(1), m.group(2)
            dim_key = jk_dim_map.get(dim_short)
            if not dim_key or dim_key not in jk_seg_ts:
                continue
            if seg_name not in jk_seg_ts[dim_key]:
                continue
            seg_data = jk_seg_ts[dim_key][seg_name]
            reason_seg = jk_seg_reasons.get(dim_key, {}).get(seg_name, {})
            ok = update_dc_slide(slide, seg_data, seg_name, dim_short,
                                 jk_seg_quarters, jk_channels_4sa, jk_ch_display,
                                 jk_display_to_raw, reason_seg,
                                 jk_main_ch, jk_threshold_ch,
                                 '잡코리아', '사람인',
                                 line_colors=JK_LINE_COLORS, bar_colors=JK_REASON_COLORS, max_q=12)
            if ok:
                print(f"  Slide {idx+1}: JK [{dim_short}] {seg_name} ✓")

    # Slide 19: 재지원 (별도 처리 — 위에서 이미 완료)
    # Slide 20-22: JK 소득 (저소득/중소득/고소득)
    print("  --- JK 소득 ---")
    run_jk_dc(range(19, 22), '소득')
    # Slide 23-34: JK 산업
    print("  --- JK 산업 ---")
    run_jk_dc(range(22, 34), '산업')
    # Slide 35-50: JK 직무
    print("  --- JK 직무 ---")
    run_jk_dc(range(34, 50), '직무')
    # Slide 51-55: JK 지역
    print("  --- JK 지역 ---")
    run_jk_dc(range(50, 55), '지역')

    # AM Double Click (Slide 57-73)
    am_dim_map = {'가계소득': '소득', '근무형태': '근무형태', '직무': '직무', '지역': '지역'}
    am_main_ch = ['알바몬', '알바천국', '당근(알바)']  # 3사 항상 표시
    am_threshold_ch = []  # AM은 조건부 채널 없음

    for idx in range(56, 73):  # Slide 57-73 (AM Double Click 전체, 지역 포함)
        slide = prs.slides[idx]
        title = ''
        for shape in slide.shapes:
            if shape.has_text_frame and 'Double Click' in shape.text_frame.text:
                title = shape.text_frame.text.strip()
                break
        if not title or '재지원' in title:
            continue

        m = re.match(r'\[(.+?)\]\s*(.+?)\s*Double Click', title)
        if not m:
            continue
        dim_short, seg_name = m.group(1), m.group(2)
        dim_key = am_dim_map.get(dim_short)
        if not dim_key or dim_key not in am_seg_ts:
            continue
        if seg_name not in am_seg_ts[dim_key]:
            continue

        seg_data = am_seg_ts[dim_key][seg_name]
        reason_seg = am_seg_reasons.get(dim_key, {}).get(seg_name, {})

        ok = update_dc_slide(slide, seg_data, seg_name, dim_short,
                             am_seg_quarters, am_channels_3sa, am_ch_display,
                             am_display_to_raw, reason_seg,
                             am_main_ch, am_threshold_ch,
                             '알바몬', '알바천국',
                             line_colors=AM_LINE_COLORS, bar_colors=AM_REASON_COLORS, max_q=12)
        if ok:
            print(f"  Slide {idx+1}: AM [{dim_short}] {seg_name} ✓")

    # ========================
    # 시사점 생성 + 텍스트 삽입 (Slide 4-6, 11-13)
    # ========================
    print(f"\n[5] Generating insights...")

    # Slide 4: JK 전체 채널
    txt = gen_insight_total(jk_total, jk_all_ch, '잡코리아', q_jk[-1], q_jk[-2])
    if set_insight_text(prs.slides[3], txt):
        print(f"  Slide 4: \"{txt[:60]}...\"")

    # Slide 5: JK 온라인 4사
    txt = gen_insight_online(jk_online, jk_4, '잡코리아', '사람인')
    txt = txt.replace('잡코리아','JK').replace('사람인','SRI').replace('원티드','WTD').replace('리멤버','RMB')
    if set_insight_text(prs.slides[4], txt):
        print(f"  Slide 5: \"{txt[:60]}...\"")

    # Slide 6: JK 인지/지원/재지원
    txt = gen_insight_iar(jk_iar, 'JK')
    if set_insight_text(prs.slides[5], txt):
        print(f"  Slide 6: \"{txt[:60]}...\"")

    # Slide 11: AM 전체 채널
    txt = gen_insight_total(am_total, am_all_ch, '알바몬', q_am[-1], q_am[-2])
    if set_insight_text(prs.slides[10], txt):
        print(f"  Slide 11: \"{txt[:60]}...\"")

    # Slide 12: AM 온라인 3사
    txt = gen_insight_online(am_online, am_3, '알바몬', '알바천국')
    txt = txt.replace('알바몬','AM').replace('알바천국','AH').replace('당근(알바)','DG')
    if set_insight_text(prs.slides[11], txt):
        print(f"  Slide 12: \"{txt[:60]}...\"")

    # Slide 13: AM 인지/지원/재지원
    txt = gen_insight_iar(am_iar, 'AM')
    if set_insight_text(prs.slides[12], txt):
        print(f"  Slide 13: \"{txt[:60]}...\"")

    # ========================
    # 전역 분기 텍스트 치환 (미처리 슬라이드 포함)
    # ========================
    print(f"\n[7] Global quarter text cleanup...")
    replaced_count = 0
    for sl_idx in range(len(prs.slides)):
        for shape in iter_all_shapes(prs.slides[sl_idx]):
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        old = run.text
                        # JK: 25Q3 → 25Q4
                        if '25Q3' in run.text:
                            run.text = run.text.replace('25Q3', '25Q4')
                        # AM: 25.3Q → 25.4Q
                        if '25.3Q' in run.text:
                            run.text = run.text.replace('25.3Q', '25.4Q')
                        if old != run.text:
                            replaced_count += 1
    print(f"  {replaced_count} text runs updated")

    # 저장
    print(f"\n[8] Saving: {OUTPUT.name}")
    prs.save(str(OUTPUT))
    print("Done!")


if __name__ == '__main__':
    main()
